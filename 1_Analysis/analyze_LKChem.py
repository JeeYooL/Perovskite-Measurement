import pandas as pd
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, simpledialog, colorchooser
import os
import glob
import io
import itertools
import sys
import pickle
import numpy as np
from scipy.signal import savgol_filter
import re
from collections import defaultdict

# --- NEW: python-pptx 라이브러리 import ---
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    messagebox.showerror("Dependency Missing", "python-pptx library is not installed.\nPlease run: python -m pip install python-pptx")
    sys.exit()
# --- END NEW ---

# --- ttkthemes 라이브러리 import ---
try:
    from ttkthemes import ThemedTk
except ImportError:
    messagebox.showerror("Dependency Missing", "ttkthemes library is not installed.\nPlease run: python -m pip install ttkthemes")
    sys.exit()

# --- Pillow, pywin32 라이브러리 import (클립보드 복사 기능) ---
try:
    from PIL import Image
except ImportError:
    messagebox.showerror("Dependency Missing", "Pillow library is not installed.\nPlease run: python -m pip install Pillow")
    sys.exit()

if sys.platform == "win32":
    try:
        import win32clipboard
    except ImportError:
        messagebox.showerror("Dependency Missing", "pywin32 library is not installed.\nPlease run: python -m pip install pywin32")
        sys.exit()

# --- 그래프 연동을 위한 matplotlib import ---
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg, NavigationToolbar2Tk

# --- 학술지 표준 폰트 및 스타일 설정 ---
plt.rcParams['font.family'] = 'sans-serif'
plt.rcParams['font.sans-serif'] = ['Arial', 'Helvetica', 'DejaVu Sans']
plt.rcParams['font.size'] = 9
plt.rcParams['axes.titlesize'] = 10
plt.rcParams['axes.labelsize'] = 9
plt.rcParams['xtick.labelsize'] = 8
plt.rcParams['ytick.labelsize'] = 8
plt.rcParams['legend.fontsize'] = 8
plt.rcParams['axes.linewidth'] = 1.0
plt.rcParams['lines.linewidth'] = 2.0
plt.rcParams['lines.markersize'] = 5.0


# --- 데이터 파싱 함수 ---
def parse_value(value_str):
    if value_str is None: return None
    value_str = value_str.strip()
    try:
        if value_str.endswith('m'): return float(value_str[:-1]) / 1000.0
        elif value_str.endswith('u'): return float(value_str[:-1]) / 1000000.0
        elif value_str.endswith('k'): return float(value_str[:-1]) * 1000.0
        else: return float(value_str)
    except (ValueError, TypeError):
        return np.nan
    except Exception:
        return np.nan

# --- Series/Shunt Resistance 계산 함수 ---
def calculate_resistances(df):
    Rs, Rsh = None, None
    try:
        V_raw = df['V'].values
        J_raw = df['J(A/cm2)'].values

        if not np.issubdtype(V_raw.dtype, np.number) or not np.issubdtype(J_raw.dtype, np.number):
            return None, None

        valid_mask = ~np.isnan(V_raw) & ~np.isnan(J_raw)
        if np.sum(valid_mask) < 10:
             return None, None
        V_raw = V_raw[valid_mask]
        J_raw = J_raw[valid_mask]

        window_length = 5
        polyorder = 3

        if len(V_raw) < window_length:
            window_length = len(V_raw) if len(V_raw) % 2 != 0 else len(V_raw) - 1
            if window_length < 3: return None, None

        J_smooth = savgol_filter(J_raw, window_length, polyorder)

        mask_rsh = np.abs(V_raw) < 0.1
        if np.sum(mask_rsh) > 2:
            coeffs_rsh = np.polyfit(V_raw[mask_rsh], J_smooth[mask_rsh], 1)
            slope_dJ_dV_at_V0 = coeffs_rsh[0]
            if slope_dJ_dV_at_V0 != 0:
                Rsh = np.abs(1.0 / slope_dJ_dV_at_V0)

        voc_index = np.argmin(np.abs(J_smooth))
        voc_value = V_raw[voc_index]
        mask_rs = (V_raw > voc_value - 0.05) & (V_raw < voc_value + 0.05)

        if np.sum(mask_rs) > 2:
            coeffs_rs = np.polyfit(V_raw[mask_rs], J_smooth[mask_rs], 1)
            slope_dJ_dV_at_Voc = coeffs_rs[0]
            if slope_dJ_dV_at_Voc != 0:
                Rs = np.abs(1.0 / slope_dJ_dV_at_Voc)

        return Rs, Rsh

    except Exception as e:
        return None, None

# --- Scan Direction 감지 함수 ---
def detect_scan_direction(file_path, df):
    file_name = os.path.basename(file_path).lower()
    if 'rev' in file_name or 'reverse' in file_name: return 'Reverse'
    elif 'fwd' in file_name or 'forward' in file_name: return 'Forward'
    try:
        V = df['V'].values
        if len(V) > 1:
            if V[0] < V[-1]: return 'Forward'
            elif V[0] > V[-1]: return 'Reverse'
    except: pass
    return 'Unknown'

# --- 전역 변수 ---
jv_data_cache = {}
original_all_devices_df = pd.DataFrame()
current_display_df = pd.DataFrame()
plotted_jv_items = {}
plotted_dist_items = {}

jv_color_cycle = itertools.cycle([mcolors.to_hex(c) for c in plt.colormaps['Dark2'].colors])
jv_marker_cycle = itertools.cycle(['o', 's', '^', 'D', 'v', 'p', '*'])
jv_linestyle_cycle = itertools.cycle(['-', '--', ':', '-.'])
dist_color_cycle = itertools.cycle([mcolors.to_hex(c) for c in plt.colormaps['Set2'].colors])

color_image_cache = {}
pce_data_by_folder = {}
current_root_folder = ""
operator_name = "" # [신규] 사용자 이름
device_structure = "" # [신규] p-i-n 또는 n-i-p
experimental_variables = {}
process_details = {}

variable_columns = ["TCO", "HTL", "Buried", "Perovskite", "Surface", "ETL", "Contact"]

VARIABLE_PRESETS = {
    "TCO": ["FTO (tec8)", "FTO (ashahi)", "ITO (china)", "ITO (우양)"],
    "HTL": ["NiOx", "Al2O3", "Me-4PACz", "MeO-2PACz", "2PACz", "PEDOT:PSS", "PTAA" ],
    "Buried": ["PEAI", "FAI", "MAI"],
    "ETL": ["C60", "PCBM", "BCP", "ALD", "SnO2"],
    "Surface": ["PDAI2", "PEAI", "OAI"],
    "Contact": ["Ag", "Au", "Cu", "ITO", "IZO"]
}

var_entry_widgets = {}
var_checkbox_vars = {}

# --- 유틸리티 함수 ---
def clean_column_names_for_ml(df):
    """
    XGBoost 호환을 위해 DataFrame의 모든 컬럼명에서 특수문자를 제거합니다.
    
    특수문자 (, ), [, ], <, > 등을 언더스코어(_)로 변환합니다.
    연속된 언더스코어는 하나로 축소하고, 앞뒤 언더스코어를 제거합니다.
    
    Args:
        df: pandas DataFrame
    
    Returns:
        pandas DataFrame with cleaned column names
    """
    cleaned_columns = []
    for col in df.columns:
        # 특수문자를 언더스코어로 변환
        cleaned = re.sub(r'[(){}\[\]<>·]', '_', str(col))
        # 공백을 언더스코어로 변환
        cleaned = cleaned.replace(' ', '_')
        # 연속된 언더스코어를 하나로 축소
        cleaned = re.sub(r'_+', '_', cleaned)
        # 앞뒤 언더스코어 제거
        cleaned = cleaned.strip('_')
        cleaned_columns.append(cleaned)
    
    df.columns = cleaned_columns
    return df

def create_color_image(color):
    if color in color_image_cache: return color_image_cache[color]
    image = tk.PhotoImage(width=16, height=16)
    image.put(color, to=(0, 0, 15, 15))
    color_image_cache[color] = image
    return image

def copy_figure_to_clipboard(fig, bbox_inches=None):
    if sys.platform != "win32": messagebox.showwarning("Unsupported OS", "Copying graphs to clipboard is only supported on Windows."); return
    try:
        with io.BytesIO() as buf:
            fig.savefig(buf, format='png', dpi=300, bbox_inches=bbox_inches)
            buf.seek(0)
            image = Image.open(buf)
            with io.BytesIO() as output:
                image.save(output, 'BMP'); data = output.getvalue()[14:]
        win32clipboard.OpenClipboard(); win32clipboard.EmptyClipboard()
        win32clipboard.SetClipboardData(win32clipboard.CF_DIB, data)
        win32clipboard.CloseClipboard(); messagebox.showinfo("Copied", "Graph has been copied to the clipboard.")
    except Exception as e:
        messagebox.showerror("Copy Error", f"Failed to copy graph to clipboard:\n{e}")

def copy_subplot_to_clipboard(index):
    ax = axs_dist.flat[index]
    extent = ax.get_tightbbox(fig_dist.canvas.get_renderer()).transformed(fig_dist.dpi_scale_trans.inverted())
    extent = extent.expanded(1.2, 1.2) 
    copy_figure_to_clipboard(fig_dist, bbox_inches=extent)

# --- Treeview Control Functions ---
def expand_all_folders():
    for item_id in file_tree.get_children():
        file_tree.item(item_id, open=True)

def collapse_all_folders():
    for item_id in file_tree.get_children():
        file_tree.item(item_id, open=False)

# --- [추가] 샘플 폴더 이름 변경을 위한 함수 ---

def rename_sample(item_id, old_sample_name):
    """샘플 폴더의 이름을 변경하는 로직을 처리합니다."""
    # 데이터가 저장된 모든 전역 변수를 가져옵니다.
    global current_display_df, original_all_devices_df, experimental_variables, process_details, pce_data_by_folder, plotted_dist_items
    
    new_sample_name = simpledialog.askstring(
        "Rename Sample", 
        f"'{old_sample_name}'의 새 이름을 입력하세요:", 
        initialvalue=old_sample_name, 
        parent=root
    )
    
    # 사용자가 취소했거나, 이름이 변경되지 않았거나, 공백을 입력한 경우
    if not new_sample_name or new_sample_name == old_sample_name or not new_sample_name.strip():
        return
        
    new_sample_name = new_sample_name.strip()

    # 이미 이름이 존재하는지 확인합니다.
    if new_sample_name in current_display_df['Sample'].unique():
        messagebox.showerror("Name Exists", f"샘플 이름 '{new_sample_name}'이(가) 이미 존재합니다. 다른 이름을 선택하세요.")
        return

    # 1. 현재 표시 중인 DataFrame의 'Sample' 열 값을 변경합니다.
    current_display_df.loc[current_display_df['Sample'] == old_sample_name, 'Sample'] = new_sample_name
    
    # 2. 원본 DataFrame의 값도 변경합니다. (필터 초기화 시 필요)
    original_all_devices_df.loc[original_all_devices_df['Sample'] == old_sample_name, 'Sample'] = new_sample_name
    
    # 3. experimental_variables 딕셔너리의 키를 변경합니다.
    if old_sample_name in experimental_variables:
        experimental_variables[new_sample_name] = experimental_variables.pop(old_sample_name)
        
    # 4. process_details 딕셔너리의 키를 변경합니다.
    if old_sample_name in process_details:
        process_details[new_sample_name] = process_details.pop(old_sample_name)
        
    # 5. 분포 그래프(dist_plot)에 사용된 딕셔너리의 'label' 값을 변경합니다.
    for item_data in plotted_dist_items.values():
        if item_data['label'] == old_sample_name:
            item_data['label'] = new_sample_name
            # 분포 그래프 목록(dist_list_tree)의 텍스트도 업데이트합니다.
            for item in dist_list_tree.get_children():
                if dist_list_tree.item(item, 'text') == old_sample_name:
                    dist_list_tree.item(item, text=new_sample_name)
                    break

    # 6. 모든 테이블과 파일 트리를 새 이름으로 새로 고칩니다.
    refresh_all_views(current_display_df)
    
    # 7. 분포 그래프의 축 라벨을 새 이름으로 다시 그립니다.
    redraw_dist_plot()

def on_file_tree_right_click(event):
    """file_tree에서 마우스 오른쪽 버튼 클릭 시 컨텍스트 메뉴를 표시합니다."""
    item_id = file_tree.identify_row(event.y)
    if not item_id:
        return # 빈 공간 클릭 시 무시

    # 클릭한 항목이 자식이 없는 최상위 항목(샘플 폴더)인지 확인합니다.
    if file_tree.parent(item_id) == '':
        file_tree.selection_set(item_id) # 클릭한 항목을 선택 상태로 만듭니다.
        old_sample_name = file_tree.item(item_id, 'text')
        
        # 팝업 메뉴 생성
        menu = tk.Menu(root, tearoff=0)
        menu.add_command(
            label=f"'{old_sample_name}' 이름 변경...", 
            command=lambda i=item_id, o=old_sample_name: rename_sample(i, o)
        )
        
        # 현재 마우스 위치에 메뉴를 띄웁니다.
        try:
            menu.tk_popup(event.x_root, event.y_root)
        finally:
            menu.grab_release()

# --- [추가 끝] ---


# --- J-V 그래프 함수들 ---
def redraw_jv_graphs():
    ax_jv.clear()
    if not plotted_jv_items:
        ax_jv.set_title("J-V Curve")
        ax_jv.set_xlabel("Voltage (V)"); ax_jv.set_ylabel("Current Density (mA/cm²)")
    else:
        for item_data in plotted_jv_items.values():
            df = item_data['df']
            label = item_data['label']
            if 'scan_dir' in item_data:
                label += f" [{item_data['scan_dir']}]"
            ax_jv.plot(df['V'], df['J(A/cm2)'] * 1000,
                       label=label,
                       color=item_data['color'],
                       marker=item_data['marker'],
                       linestyle=item_data['linestyle'],
                       markersize=5,
                       linewidth=2
                       )
        ax_jv.legend(); ax_jv.set_title("J-V Curve")
        ax_jv.set_xlabel("Voltage (V)"); ax_jv.set_ylabel("Current Density (mA/cm²)")
    ax_jv.grid(True, linestyle='--', alpha=0.6); ax_jv.axhline(0, color='gray', linewidth=0.5); ax_jv.axvline(0, color='gray', linewidth=0.5)
    fig_jv.tight_layout(pad=2.0); jv_canvas_widget.draw()

def add_selected_to_jv_graph():
    global current_display_df
    selected_items = file_tree.selection()
    if not selected_items: messagebox.showwarning("Selection Needed", "Please select file(s) from the list to add to the J-V plot."); return

    scan_filter = scan_filter_jv_var.get()

    for item_id in selected_items:
        if file_tree.parent(item_id):
            file_path = file_tree.item(item_id, 'values')[0]
            if file_path not in [d['path'] for d in plotted_jv_items.values()]:
                df = jv_data_cache.get(file_path)
                if df is not None:
                    if file_path not in current_display_df['FullPath'].values: continue
                    file_info = current_display_df[current_display_df['FullPath'] == file_path].iloc[0]
                    scan_dir = file_info['Scan Direction']

                    if scan_filter != 'All' and scan_dir != scan_filter: continue

                    label = os.path.basename(file_path)
                    color = next(jv_color_cycle)
                    marker = next(jv_marker_cycle)
                    linestyle = next(jv_linestyle_cycle)

                    color_img = create_color_image(color)
                    new_item_id = plotted_list_tree.insert('', 'end', text=label, image=color_img)

                    plotted_jv_items[new_item_id] = {
                        'path': file_path, 'label': label, 'color': color,
                        'df': df, 'image': color_img, 'scan_dir': scan_dir,
                        'marker': marker, 'linestyle': linestyle
                    }
    redraw_jv_graphs(); auto_scale_jv_axes(smart_scale=True)

def remove_selected_from_jv_graph():
    selected_items = plotted_list_tree.selection()
    if not selected_items: messagebox.showwarning("Selection Needed", "Please select item(s) from the 'Plotted J-V Curves' list to remove."); return
    for item_id in selected_items:
        if item_id in plotted_jv_items: del plotted_jv_items[item_id]
        plotted_list_tree.delete(item_id)
    redraw_jv_graphs()

def clear_jv_plot():
    plotted_jv_items.clear(); plotted_list_tree.delete(*plotted_list_tree.get_children())
    global jv_color_cycle, jv_marker_cycle, jv_linestyle_cycle
    jv_color_cycle = itertools.cycle([mcolors.to_hex(c) for c in plt.colormaps['Dark2'].colors])
    jv_marker_cycle = itertools.cycle(['o', 's', '^', 'D', 'v', 'p', '*'])
    jv_linestyle_cycle = itertools.cycle(['-', '--', ':', '-.'])
    redraw_jv_graphs()

def on_plotted_item_double_click(event, item_dict, tree_widget):
    selected_item_id = tree_widget.focus()
    if not selected_item_id: return
    item_data = item_dict.get(selected_item_id)
    if not item_data: return

    new_label = simpledialog.askstring("Change Label", "Enter new graph label:", initialvalue=item_data['label'])
    if new_label: item_data['label'] = new_label
    color_code = colorchooser.askcolor(title="Choose color", initialcolor=item_data['color'])
    if color_code and color_code[1]: item_data['color'] = color_code[1]

    new_color_img = create_color_image(item_data['color'])
    item_data['image'] = new_color_img
    tree_widget.item(selected_item_id, text=item_data['label'], image=new_color_img)

    if tree_widget == plotted_list_tree: redraw_jv_graphs()
    elif tree_widget == dist_list_tree: redraw_dist_plot()

def update_jv_axis_limits():
    try:
        xmin = float(jv_x_min_var.get()) if jv_x_min_var.get() else None; xmax = float(jv_x_max_var.get()) if jv_x_max_var.get() else None
        ymin = float(jv_y_min_var.get()) if jv_y_min_var.get() else None; ymax = float(jv_y_max_var.get()) if jv_y_max_var.get() else None
        ax_jv.set_xlim(xmin, xmax); ax_jv.set_ylim(ymin, ymax); jv_canvas_widget.draw()
    except (ValueError, TypeError): messagebox.showerror("Input Error", "J-V axis range must be numbers.")

def auto_scale_jv_axes(smart_scale=False):
    if smart_scale and plotted_jv_items:
        max_x, max_y = -np.inf, -np.inf
        for item in plotted_jv_items.values():
            max_x = max(max_x, item['df']['V'].max()); max_y = max(max_y, (item['df']['J(A/cm2)'] * 1000).max())
        ax_jv.set_xlim(0, max_x * 1.1 if max_x > 0 else 1.0); ax_jv.set_ylim(0, max_y * 1.1 if max_y > 0 else 25.0)
    else: ax_jv.autoscale_view()
    jv_x_min_var.set(""); jv_x_max_var.set(""); jv_y_min_var.set(""); jv_y_max_var.set(""); jv_canvas_widget.draw()

def plot_best_pces():
    global current_display_df
    if current_display_df.empty: messagebox.showwarning("No Data", "Please load and filter data first."); return

    clear_jv_plot()

    scan_filter = scan_filter_jv_var.get()
    df_filtered = current_display_df.copy()
    if scan_filter != 'All':
        df_filtered = df_filtered[df_filtered['Scan Direction'] == scan_filter]

    if df_filtered.empty:
        messagebox.showwarning("No Data", f"No {scan_filter} scan data available.")
        return

    best_devices_idx = df_filtered.groupby('Sample')['PCE (%)'].idxmax()
    best_devices_df = df_filtered.loc[best_devices_idx]

    for index, row in best_devices_df.iterrows():
        file_path = row['FullPath']
        sample_name = row['Sample']
        scan_dir = row['Scan Direction']
        df = jv_data_cache.get(file_path)

        if df is not None:
            label = sample_name

            color = next(jv_color_cycle)
            marker = next(jv_marker_cycle)
            linestyle = next(jv_linestyle_cycle)

            color_img = create_color_image(color)
            new_item_id = plotted_list_tree.insert('', 'end', text=label, image=color_img)

            plotted_jv_items[new_item_id] = {
                'path': file_path, 'label': label, 'color': color,
                'df': df, 'image': color_img, 'scan_dir': scan_dir,
                'marker': marker, 'linestyle': linestyle
            }

    redraw_jv_graphs(); auto_scale_jv_axes(smart_scale=True)


# --- [추가] J-V 원본 데이터 복사 함수 ---
def copy_jv_raw_data():
    """현재 그려진 모든 J-V 곡선의 원본 데이터를 (V, J(mA/cm2)) 탭으로 구분된 형식으로 클립보드에 복사합니다."""
    
    # 1. 그래프에 그려진 아이템이 있는지 확인
    if not plotted_jv_items:
        messagebox.showwarning("데이터 없음", "먼저 J-V 곡선을 그래프에 추가하세요.", parent=root)
        return
    
    try:
        all_dfs_to_concat = []
        
        # 2. "Plotted J-V Curves" 목록의 순서대로 데이터를 가져옵니다.
        item_ids_in_order = plotted_list_tree.get_children()
        
        if not item_ids_in_order:
             # 목록이 비어있으면(이론상으론 없어야 함) 딕셔너리 순서로 대체
             item_ids_in_order = plotted_jv_items.keys()

        for item_id in item_ids_in_order:
            if item_id not in plotted_jv_items:
                continue
                
            data = plotted_jv_items[item_id]
            label = data['label'] # 예: "Device A"
            df = data['df'].copy()
            
            # 3. OriginPro에서 사용할 고유한 컬럼 이름을 생성합니다.
            v_col_name = f"V ({label})"       # 예: "V (Device A)"
            j_col_name = f"J_mA_cm2 ({label})" # 예: "J_mA_cm2 (Device A)"
            
            # 4. V 컬럼과 mA/cm²로 변환된 J 컬럼으로 새 DataFrame을 생성합니다.
            # reset_index(drop=True)는 서로 다른 길이의 데이터를 합칠 때 중요합니다.
            df_to_add = pd.DataFrame({
                v_col_name: df['V'],
                j_col_name: df['J(A/cm2)'] * 1000.0
            }).reset_index(drop=True)
            
            all_dfs_to_concat.append(df_to_add)

        if not all_dfs_to_concat:
            messagebox.showwarning("데이터 없음", "플롯된 항목의 데이터를 찾을 수 없습니다.", parent=root)
            return
        # 5. 모든 DataFrame을 가로(축=1)로 합칩니다. (V1, J1, V2, J2, ...)
        combined_df = pd.concat(all_dfs_to_concat, axis=1)
        
        # 6. DataFrame을 탭(tab)으로 구분된(sep='\t') 문자열로 변환합니다.
        #    정밀도를 위해 과학적 표기법(e)을 사용합니다.
        output_stream = io.StringIO()
        combined_df.to_csv(output_stream, sep='\t', index=False, float_format='%.6e')
        tsv_data = output_stream.getvalue()
        
        # 7. 클립보드에 복사합니다.
        root.clipboard_clear()
        root.clipboard_append(tsv_data)
        messagebox.showinfo("복사 완료", 
                            f"{len(all_dfs_to_concat)}개의 J-V 곡선 원본 데이터가 클립보드에 복사되었습니다.\n(탭으로 구분됨)", 
                            parent=root)

    except Exception as e:
        messagebox.showerror("복사 오류", f"원본 데이터를 복사하는 중 오류가 발생했습니다:\n{e}", parent=root)
# --- [추가 끝] ---


# --- 분포 그래프 함수들 ---
def redraw_dist_plot():
    for i, ax in enumerate(axs_dist.flat): ax.clear()

    scan_filter = scan_filter_dist_var.get()

    if not plotted_dist_items:
        param_names = ['Voc (V)', 'Jsc (mA/cm²)', 'FF (%)', 'PCE (%)']
        for i, ax in enumerate(axs_dist.flat):
            ax.set_title(param_names[i]); ax.grid(True, linestyle='--', alpha=0.6)
    else:
        data_to_plot = {'Voc (V)': [], 'Jsc (mA/cm2)': [], 'FF (%)': [], 'PCE (%)': []}
        labels, colors = [], []
        for item_data in plotted_dist_items.values():
            folder_name = item_data['label']
            df = pce_data_by_folder.get(folder_name)
            if df is not None and not df.empty:
                if scan_filter != 'All':
                    df = df[df['Scan Direction'] == scan_filter]

                if df.empty: continue

                labels.append(folder_name); colors.append(item_data['color'])
                for param in data_to_plot.keys(): data_to_plot[param].append(df[param].dropna())

        param_keys = list(data_to_plot.keys())
        for i, ax in enumerate(axs_dist.flat):
            param = param_keys[i]
            param_data = data_to_plot[param]
            if not any(len(d) > 0 for d in param_data): continue

            bp = ax.boxplot(param_data, patch_artist=True, tick_labels=labels, showfliers=False)

            for i_cat, cat_data in enumerate(param_data):
                x_coords = np.random.normal(i_cat + 1, 0.03, size=len(cat_data))
                ax.scatter(x_coords, cat_data, s=20, color=colors[i_cat], edgecolor='black', linewidths=0.5, alpha=0.7, zorder=10)

            ax.set_title(param); ax.tick_params(axis='x', rotation=15); ax.grid(True, linestyle='--', alpha=0.6)
            for patch, color in zip(bp['boxes'], colors): patch.set_facecolor(color)
    fig_dist.tight_layout(pad=2.0); dist_canvas_widget.draw()

def change_dist_layout():
    """[수정됨] Distribution 그래프의 레이아웃과 컨트롤 프레임의 위치,
       그리고 컨트롤 프레임 '내부의' 배치까지 동적으로 변경합니다."""
    
    # [수정] 컨트롤의 자식 프레임 3개를 global로 가져옵니다.
    global fig_dist, axs_dist, dist_canvas_widget, dist_toolbar, \
           dist_plot_container_frame, dist_toolbar_frame, dist_controls_area_frame, \
           dist_list_frame, dist_axis_manage_frame, dist_export_frame
    
    layout = dist_layout_var.get()
    
    # 1. 기존 Matplotlib Figure 객체 닫기
    plt.close(fig_dist)
    
    if layout == "4x1":
        # [수정] 10인치 -> 8인치로 대폭 축소 (세로가 긴 형태 유지)
        fig_dist, axs_dist = plt.subplots(4, 1, facecolor='white', figsize=(5, 8))
    elif layout == "1x4":
        # [수정] 16인치 -> 12인치, 4인치 -> 3.5인치
        fig_dist, axs_dist = plt.subplots(1, 4, facecolor='white', figsize=(12, 3.5))
    else: # 2x2 (default)
        # [수정] 8x7 -> 7x5
        fig_dist, axs_dist = plt.subplots(2, 2, facecolor='white', figsize=(7, 5)) 
    # ... (생략) ...
    
    # 3. 기존 Tkinter 위젯(Canvas) 파괴
    for widget in dist_plot_container_frame.winfo_children():
        if widget != dist_toolbar_frame: # 툴바 '프레임'은 남겨둠
            widget.destroy()

    # 4. 툴바 프레임(dist_toolbar_frame)의 내용물(버튼, 콤보박스 등)도 모두 비웁니다.
    for widget in dist_toolbar_frame.winfo_children():
        widget.destroy()

    # 5. 새 Canvas 위젯을 왼쪽 컨테이너(dist_plot_container_frame)에 다시 생성
    dist_canvas_widget = FigureCanvasTkAgg(fig_dist, master=dist_plot_container_frame)
    dist_canvas_widget.get_tk_widget().pack(side=tk.TOP, fill=tk.BOTH, expand=True)
    
    # 6. 새 Toolbar 객체를 툴바 프레임(dist_toolbar_frame)에 다시 생성
    dist_toolbar = NavigationToolbar2Tk(dist_canvas_widget, dist_toolbar_frame)
    dist_toolbar.update()

    # 7. 툴바 프레임에 있던 커스텀 버튼들을 모두 다시 생성합니다.
    ttk.Button(dist_toolbar_frame, text="Copy All", command=lambda: copy_figure_to_clipboard(fig_dist)).pack(side=tk.RIGHT, padx=5)
    
    dist_layout_var.set(layout) 
    ttk.Label(dist_toolbar_frame, text="Layout:").pack(side=tk.RIGHT, padx=(5,2))
    layout_combo = ttk.Combobox(dist_toolbar_frame, textvariable=dist_layout_var, values=['2x2', '1x4', '4x1'], state='readonly', width=6)
    layout_combo.pack(side=tk.RIGHT, padx=2)
    layout_combo.bind('<<ComboboxSelected>>', lambda e: change_dist_layout())

    current_scan_filter = scan_filter_dist_var.get() 
    ttk.Label(dist_toolbar_frame, text="Scan:").pack(side=tk.RIGHT, padx=(5,2))
    scan_combo_dist = ttk.Combobox(dist_toolbar_frame, textvariable=scan_filter_dist_var, values=['All', 'Reverse', 'Forward'], state='readonly', width=10)
    scan_combo_dist.pack(side=tk.RIGHT, padx=2)
    scan_combo_dist.set(current_scan_filter) 
    scan_combo_dist.bind('<<ComboboxSelected>>', lambda e: redraw_dist_plot())
    
    
    # 8. [핵심] 레이아웃에 따라 마스터 프레임과 그 자식들의 위치를 동적으로 재배치
    
    # 8-1. 메인 프레임 2개를 뗀다
    dist_plot_container_frame.pack_forget()
    dist_controls_area_frame.pack_forget()
    
    # 8-2. 컨트롤 프레임 내부의 자식 3개도 모두 뗀다
    dist_list_frame.pack_forget()
    dist_axis_manage_frame.pack_forget()
    dist_export_frame.pack_forget()

    if layout == '4x1':
        # --- 4x1 배치 ---
        # 8-3. 메인 프레임 배치: 그래프(LEFT), 컨트롤(RIGHT)
        dist_plot_container_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=5, pady=5)
        dist_controls_area_frame.pack(side=tk.RIGHT, fill=tk.Y, expand=False, padx=(0, 5), pady=5)
        
        # 8-4. 컨트롤 자식 3개 배치: 세로(TOP)로 쌓기
        dist_list_frame.pack(side=tk.TOP, fill=tk.X, expand=False, pady=(0, 5))
        dist_axis_manage_frame.pack(side=tk.TOP, fill=tk.X, expand=False, pady=5)
        dist_export_frame.pack(side=tk.TOP, fill=tk.X, expand=False, pady=5)
    else:
        # --- 2x2 또는 1x4 배치 ---
        # 8-3. 메인 프레임 배치: 그래프(TOP), 컨트롤(BOTTOM)
        dist_plot_container_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=5, pady=5)
        dist_controls_area_frame.pack(side=tk.BOTTOM, fill=tk.X, expand=False, padx=5, pady=5)
        
        # 8-4. 컨트롤 자식 3개 배치: 가로(LEFT) 3분할 (이미지 참고)
        dist_list_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, pady=0, padx=(0,5))
        dist_axis_manage_frame.pack(side=tk.LEFT, fill=tk.Y, expand=False, pady=0, padx=5)
        dist_export_frame.pack(side=tk.LEFT, fill=tk.Y, expand=False, pady=0, padx=(5,0))

    
    # 9. 그래프 다시 그리기
    redraw_dist_plot()

def add_selected_folder_to_dist_plot():
    selected_items = file_tree.selection()
    if not selected_items: messagebox.showwarning("Selection Needed", "Please select sample folder(s) from the list to add to the distribution plot."); return
    for item_id in selected_items:
        if not file_tree.parent(item_id):
            folder_name = file_tree.item(item_id, 'text')
            if folder_name not in [d['label'] for d in plotted_dist_items.values()]:
                color = next(dist_color_cycle)
                color_img = create_color_image(color)
                new_item_id = dist_list_tree.insert('', 'end', text=folder_name, image=color_img)
                plotted_dist_items[new_item_id] = {'label': folder_name, 'color': color, 'image': color_img}
    redraw_dist_plot()

def remove_selected_from_dist_graph():
    selected_items = dist_list_tree.selection()
    if not selected_items: messagebox.showwarning("Selection Needed", "Please select item(s) from the 'Plotted Distributions' list to remove."); return
    for item_id in selected_items:
        if item_id in plotted_dist_items: del plotted_dist_items[item_id]
        dist_list_tree.delete(item_id)
    redraw_dist_plot()

def clear_dist_plot():
    plotted_dist_items.clear(); dist_list_tree.delete(*dist_list_tree.get_children())
    global dist_color_cycle
    dist_color_cycle = itertools.cycle([mcolors.to_hex(c) for c in plt.colormaps['Set2'].colors])
    redraw_dist_plot()

def update_dist_axis_limits():
    vars_list = [(axs_dist[0,0], dist_voc_min_var, dist_voc_max_var), (axs_dist[0,1], dist_jsc_min_var, dist_jsc_max_var),
                 (axs_dist[1,0], dist_ff_min_var, dist_ff_max_var), (axs_dist[1,1], dist_pce_min_var, dist_pce_max_var)]
    try:
        for ax, min_var, max_var in vars_list:
            ymin = float(min_var.get()) if min_var.get() else None; ymax = float(max_var.get()) if max_var.get() else None
            ax.set_ylim(ymin, ymax)
        dist_canvas_widget.draw()
    except (ValueError, TypeError): messagebox.showerror("Input Error", "Distribution axis range must be numbers.")

def auto_scale_dist_axes():
    for ax in axs_dist.flat: ax.autoscale_view()
    for var in [dist_voc_min_var, dist_voc_max_var, dist_jsc_min_var, dist_jsc_max_var,
                dist_ff_min_var, dist_ff_max_var, dist_pce_min_var, dist_pce_max_var]: var.set("")
    dist_canvas_widget.draw()

# --- [신규] 폴더 로드 시 사용자 정보 입력을 위한 팝업 클래스 ---
class LoadInfoDialog(simpledialog.Dialog):
    """폴더 로드 시 사용자 이름과 소자 구조를 입력받는 팝업."""
    
    def body(self, master):
        self.entries = {}

        ttk.Label(master, text="Operator Name:").grid(row=0, column=0, sticky='w', padx=5, pady=5)
        self.entries['operator'] = ttk.Entry(master, width=25)
        self.entries['operator'].grid(row=0, column=1, padx=5, pady=5)
        
        ttk.Label(master, text="Device Structure:").grid(row=1, column=0, sticky='w', padx=5, pady=5)
        
        self.structure_var = tk.StringVar(value="p-i-n") # 기본값
        structure_combo = ttk.Combobox(master, 
                                       textvariable=self.structure_var, 
                                       values=["p-i-n", "n-i-p"], 
                                       state='readonly',
                                       width=23)
        structure_combo.grid(row=1, column=1, padx=5, pady=5)
        
        self.entries['structure'] = self.structure_var 

        return self.entries['operator'] # 첫 번째 입력창에 포커스

    def apply(self):
        """OK 버튼 클릭 시 호출됨"""
        self.result = {
            'operator': self.entries['operator'].get().strip(),
            'structure': self.entries['structure'].get()
        }

# --- 데이터 로딩, 필터링, UI 업데이트 함수 ---
def refresh_all_views(df_to_display):
    global pce_data_by_folder, current_display_df
    current_display_df = df_to_display
    pce_data_by_folder.clear()
    for tree in [file_tree, pce_table_all, pce_table_best, pce_table_stats]: tree.delete(*tree.get_children())

    active_vars = [col for col in variable_columns if any(sample in experimental_variables and col in experimental_variables[sample] and experimental_variables[sample][col] for sample in df_to_display['Sample'].unique())]

    all_table_cols = ['File', 'Scan', 'Voc (V)', 'Jsc (mA/cm2)', 'FF (%)', 'PCE (%)', 'Rs (Ω·cm²)', 'Rsh (Ω·cm²)'] + active_vars
    best_table_cols = ['Sample', 'File', 'Scan', 'Voc (V)', 'Jsc (mA/cm2)', 'FF (%)', 'PCE (%)', 'Rs (Ω·cm²)', 'Rsh (Ω·cm²)'] + active_vars
    stats_table_cols = ['Sample', 'Count', 'Voc (V)', 'Jsc (mA/cm2)', 'FF (%)', 'PCE (%)', 'Rs (Ω·cm²)', 'Rsh (Ω·cm²)'] + active_vars

    for table, cols in [(pce_table_all, all_table_cols), (pce_table_best, best_table_cols), (pce_table_stats, stats_table_cols)]:
        table['columns'] = cols
        for col in cols:
            table.heading(col, text=col, anchor='center'); table.column(col, anchor='center', width=120)
        table.column('#0', width=0, stretch=tk.NO); table.heading('#0', text='', anchor='w')

        if 'File' in cols:
            table.column('File', width=300, anchor='w')
        if 'Sample' in cols:
            table.column('Sample', width=150, anchor='w')
        if 'Scan' in cols:
            table.column('Scan', width=80, anchor='center')

    pce_table_stats.column('Voc (V)', width=150); pce_table_stats.column('Jsc (mA/cm2)', width=150)
    pce_table_stats.column('FF (%)', width=150); pce_table_stats.column('PCE (%)', width=150)
    pce_table_stats.column('Rs (Ω·cm²)', width=150); pce_table_stats.column('Rsh (Ω·cm²)', width=150)


    folder_nodes = {}
    for index, row in current_display_df.iterrows():
        sample_folder = row['Sample']; file_name = row['File']; full_path = row['FullPath']
        scan_dir = row.get('Scan Direction', 'Unknown')
        display_name = f"{file_name} [{scan_dir}]"

        if sample_folder not in folder_nodes:
            parent_node = file_tree.insert('', 'end', text=sample_folder, open=False)
            folder_nodes[sample_folder] = parent_node
        file_tree.insert(folder_nodes[sample_folder], 'end', text=display_name, values=(full_path,))

        var_values = [row.get(col, "") for col in active_vars]
        pce_table_all.insert('', 'end', values=(
            f"{sample_folder}/{file_name}", scan_dir,
            f"{row['Voc (V)']:.4f}" if pd.notna(row['Voc (V)']) else "N/A", f"{row['Jsc (mA/cm2)']:.2f}" if pd.notna(row['Jsc (mA/cm2)']) else "N/A",
            f"{row['FF (%)']:.2f}" if pd.notna(row['FF (%)']) else "N/A", f"{row['PCE (%)']:.2f}" if pd.notna(row['PCE (%)']) else "N/A",
            f"{row['Rs (Ω·cm²)']:.2f}" if pd.notna(row['Rs (Ω·cm²)']) else "N/A", f"{row['Rsh (Ω·cm²)']:.1f}" if pd.notna(row['Rsh (Ω·cm²)']) else "N/A",
            *var_values
        ))

    for sample_folder, group_df in current_display_df.groupby('Sample'):
        pce_data_by_folder[sample_folder] = group_df
        df_numeric = group_df.dropna(subset=['PCE (%)'])
        if df_numeric.empty: continue
        best_device = df_numeric.loc[df_numeric['PCE (%)'].idxmax()]

        var_values = [best_device.get(col, "") for col in active_vars]
        pce_table_best.insert('', 'end', values=(
            sample_folder, best_device['File'], best_device.get('Scan Direction', 'Unknown'),
            f"{best_device['Voc (V)']:.4f}", f"{best_device['Jsc (mA/cm2)']:.2f}",
            f"{best_device['FF (%)']:.2f}", f"{best_device['PCE (%)']:.2f}",
            f"{best_device['Rs (Ω·cm²)']:.2f}" if pd.notna(best_device.get('Rs (Ω·cm²)', np.nan)) else "N/A",
            f"{best_device['Rsh (Ω·cm²)']:.1f}" if pd.notna(best_device.get('Rsh (Ω·cm²)', np.nan)) else "N/A",
            *var_values
        ))

        stats = df_numeric.describe(); mean = stats.loc['mean']; std = stats.loc['std']
        var_values = [group_df.iloc[0].get(col, "") for col in active_vars]
        pce_table_stats.insert('', 'end', values=(
            sample_folder, f"{len(df_numeric)}",
            f"{mean.get('Voc (V)', 0):.4f} ± {std.get('Voc (V)', 0):.4f}", f"{mean.get('Jsc (mA/cm2)', 0):.2f} ± {std.get('Jsc (mA/cm2)', 0):.2f}",
            f"{mean.get('FF (%)', 0):.2f} ± {std.get('FF (%)', 0):.2f}", f"{mean.get('PCE (%)', 0):.2f} ± {std.get('PCE (%)', 0):.2f}",
            f"{mean.get('Rs (Ω·cm²)', 0):.2f} ± {std.get('Rs (Ω·cm²)', 0):.2f}" if 'Rs (Ω·cm²)' in mean else "N/A",
            f"{mean.get('Rsh (Ω·cm²)', 0):.1f} ± {std.get('Rsh (Ω·cm²)', 0):.1f}" if 'Rsh (Ω·cm²)' in mean else "N/A",
            *var_values
        ))

    clear_jv_plot()
    redraw_dist_plot()

    if len(df_to_display) < len(original_all_devices_df) or 'merged' in df_to_display.columns:
        messagebox.showinfo("Update Complete", f"{len(current_display_df)} devices are now displayed.")


def apply_filter():
    global current_display_df
    if original_all_devices_df.empty: messagebox.showwarning("No Data", "Please load a folder first."); return
    try:
        voc_min = float(filter_voc_min.get() or -np.inf); voc_max = float(filter_voc_max.get() or np.inf)
        jsc_min = float(filter_jsc_min.get() or -np.inf); jsc_max = float(filter_jsc_max.get() or np.inf)
        ff_min = float(filter_ff_min.get() or -np.inf); ff_max = float(filter_ff_max.get() or np.inf)
        pce_min = float(filter_pce_min.get() or -np.inf); pce_max = float(filter_pce_max.get() or np.inf)
    except (ValueError, TypeError): messagebox.showerror("Input Error", "Filter values must be numbers."); return

    # Start filtering from the original data, but preserve merges
    temp_df_for_filter = original_all_devices_df.copy()
    
    # Apply current merge state if any merges happened before filtering
    merge_map = pd.Series(current_display_df.Sample.values, index=current_display_df.FullPath).to_dict()
    temp_df_for_filter['Sample'] = temp_df_for_filter['FullPath'].map(merge_map).fillna(temp_df_for_filter['Sample'])


    filtered_df = temp_df_for_filter[
        (temp_df_for_filter['Voc (V)'].between(voc_min, voc_max, inclusive='both')) &
        (temp_df_for_filter['Jsc (mA/cm2)'].between(jsc_min, jsc_max, inclusive='both')) &
        (temp_df_for_filter['FF (%)'].between(ff_min, ff_max, inclusive='both')) &
        (temp_df_for_filter['PCE (%)'].between(pce_min, pce_max, inclusive='both'))
    ]

    refresh_all_views(filtered_df)
    messagebox.showinfo("Filter Applied", f"{len(filtered_df)} devices match the filter criteria.")


def reset_filter():
    global current_display_df
    if original_all_devices_df.empty: return
    for var in [filter_voc_min, filter_voc_max, filter_jsc_min, filter_jsc_max, filter_ff_min, filter_ff_max, filter_pce_min, filter_pce_max]:
        var.set("")

    # Reset view to the original data, but maintain current merge state
    current_display_df = original_all_devices_df.copy()
    # Re-apply merge state from the currently displayed (potentially merged) df before reset
    merge_map = {}
    if 'Sample' in current_display_df.columns: # Check if merge happened
         merge_map = pd.Series(current_display_df.Sample.values, index=current_display_df.FullPath).to_dict()

    # Apply the map to the freshly copied original data
    current_display_df = original_all_devices_df.copy()
    if merge_map :
        current_display_df['Sample'] = current_display_df['FullPath'].map(merge_map).fillna(current_display_df['Sample'])


    refresh_all_views(current_display_df)

def filter_best_pce_per_folder():
    """
    각 샘플 폴더(Sample) 별로 PCE가 가장 높은 단 하나의 데이터만 남깁니다.
    """
    global current_display_df
    if current_display_df.empty:
        messagebox.showwarning("No Data", "Please load data first.")
        return

    try:
        # 1. 'Sample' 컬럼을 기준으로 그룹을 묶고, 각 그룹에서 'PCE (%)'가 최대인 행의 인덱스(idxmax)를 찾습니다.
        # 이렇게 하면 각 폴더당 정확히 1개의 행(최고 효율)만 선택됩니다.
        best_indices = current_display_df.groupby('Sample')['PCE (%)'].idxmax()
        
        # 2. 찾은 인덱스에 해당하는 데이터만 추출합니다.
        filtered_df = current_display_df.loc[best_indices]
        
        # 3. 화면을 갱신합니다.
        refresh_all_views(filtered_df)
        
        # 결과 메시지
        messagebox.showinfo("Best per Folder Applied", 
                            f"각 폴더별 최고 효율 데이터만 남겼습니다.\n\n"
                            f"변경 전: {len(current_display_df)} 개\n"
                            f"변경 후: {len(filtered_df)} 개")

    except Exception as e:
        messagebox.showerror("Error", f"Failed to filter best files:\n{e}")

# --- [추가] 새로운 CSV 포맷 파싱 헬퍼 함수 ---
def parse_new_csv_format(file_content):
    """
    새로운 CSV 포맷(Timestamps 헤더 및 하단 Repeat 파라미터)을 파싱하여
    DataFrame(J-V data)과 Parameters(dict)를 반환합니다.
    """
    try:
        # 데이터 영역과 파라미터 영역의 시작 줄 찾기
        header_idx = -1
        param_header_idx = -1
        
        # 리스트로 변환 (이미 리스트일 수 있음)
        lines = [line.strip() for line in file_content]
        
        for i, line in enumerate(lines):
            if line.startswith("Timestamps (sec)"):
                header_idx = i
            if line.startswith("Repeat,V oc (V)"):
                param_header_idx = i
        
        if header_idx == -1:
            return None, {} # 유효한 CSV 형식이 아님

        # 1. J-V 데이터 파싱
        # 파라미터 섹션 전까지, 혹은 빈 줄 전까지가 데이터
        end_data_idx = param_header_idx if param_header_idx != -1 else len(lines)
        
        # 데이터 영역 문자열 재구성 (헤더 포함)
        data_lines = lines[header_idx:end_data_idx]
        # 빈 줄 제거 (데이터 끝부분)
        data_lines = [l for l in data_lines if l.strip()]
        
        data_str = "\n".join(data_lines)
        df = pd.read_csv(io.StringIO(data_str))
        
        # 컬럼 이름 변경 및 단위 변환 (기존 코드 호환성 유지)
        # New: Voltage (V) -> Old: V
        # New: Current Density (mA/cm^2) -> Old: J(A/cm2) (단위 변환 필요)
        
        rename_map = {
            'Voltage (V)': 'V',
            'Current Density (mA/cm^2)': 'J_mA_raw' # 임시 이름
        }
        df.rename(columns=rename_map, inplace=True)
        
        # J(A/cm2) 계산: mA/cm2 / 1000
        if 'J_mA_raw' in df.columns:
            df['J(A/cm2)'] = df['J_mA_raw'] / 1000.0
        else:
            # 만약 컬럼명이 다르면 첫번째, 세번째 컬럼 추정 등으로 fallback 가능
            pass

        # 2. 파라미터 파싱
        parameters = {}
        if param_header_idx != -1 and param_header_idx + 1 < len(lines):
            # 헤더: Repeat,V oc (V),J sc (mA/cm^2),Fill Factor,Efficiency (%)
            # 데이터: 1(R),1.03E+0,24.16E+0,744.69E-3,18.54E+0
            
            param_headers = lines[param_header_idx].split(',')
            param_values = lines[param_header_idx + 1].split(',')
            
            if len(param_headers) == len(param_values):
                raw_params = dict(zip(param_headers, param_values))
                
                # 기존 코드의 키 이름으로 매핑
                # Voc
                if 'V oc (V)' in raw_params:
                    parameters['Voc (V)'] = raw_params['V oc (V)']
                
                # Jsc (단위 주의: 파일은 mA/cm^2)
                if 'J sc (mA/cm^2)' in raw_params:
                    # 기존 코드 parse_value는 단위를 처리하지만, 여기선 숫자만 넘김
                    # 단위 변환은 load_and_process_folder 메인 루프에서 처리됨
                    parameters['Jsc(A/cm2)'] = str(parse_value(raw_params['J sc (mA/cm^2)']) / 1000.0)

                # Fill Factor
                # 파일 값: 744.69E-3 (= 0.74469) -> UI 기대값: 74.47 (%)
                if 'Fill Factor' in raw_params:
                    ff_val = parse_value(raw_params['Fill Factor'])
                    if ff_val < 1.0: # 비율로 저장된 경우 %로 변환
                        ff_val *= 100.0
                    parameters['Fill factor (%)'] = str(ff_val)

                # Efficiency
                if 'Efficiency (%)' in raw_params:
                    parameters['Efficiency (%)'] = raw_params['Efficiency (%)']
                
                # Scan Direction 추출 (Repeat 컬럼 이용)
                # 예: "1(R)" -> Reverse, "1(F)" -> Forward
                if 'Repeat' in raw_params:
                    rep_val = raw_params['Repeat']
                    if '(R)' in rep_val:
                        parameters['Scan Direction'] = 'Reverse'
                    elif '(F)' in rep_val:
                        parameters['Scan Direction'] = 'Forward'

        return df, parameters

    except Exception as e:
        print(f"CSV Parsing Error: {e}")
        return None, {}

# --- [수정] load_and_process_folder 함수 ---
def load_and_process_folder(folder_path=None):
    # 전역 변수 호출
    global original_all_devices_df, current_display_df, current_root_folder, \
           experimental_variables, jv_data_cache, \
           operator_name, device_structure
    
    # 1. 정보 입력 팝업 (신규 로드 시에만)
    is_fresh_load = folder_path is None
    if is_fresh_load:
        folder_path = filedialog.askdirectory(title="Select Date Folder (e.g., 20251021)")
        if not folder_path: return
        
        # 사용자 정보 입력 팝업
        dialog = LoadInfoDialog(root, title="Session Info")
        if dialog.result:
            operator_name = dialog.result['operator']
            device_structure = dialog.result['structure']
        else:
            operator_name = ""
            device_structure = ""

    # 데이터 초기화
    current_root_folder = folder_path
    experimental_variables = {}
    jv_data_cache = {}
    original_all_devices_df = pd.DataFrame()
    current_display_df = pd.DataFrame()
    clear_jv_plot()
    clear_dist_plot()
    reset_filter()

    # 파일 검색 (.txt 및 .csv)
    txt_files = glob.glob(os.path.join(folder_path, '**', '*.txt'), recursive=True)
    csv_files = glob.glob(os.path.join(folder_path, '**', '*.csv'), recursive=True)
    file_list = sorted(txt_files + csv_files)

    if not file_list: messagebox.showinfo("Info", "No .txt or .csv files found."); return
    all_data = []
    
    # 로딩 프로그레스 바
    progress_win = tk.Toplevel(root)
    progress_win.title("Loading Data")
    progress_win.geometry("400x100")
    progress_win.transient(root)
    progress_win.grab_set()
    ttk.Label(progress_win, text="Loading and processing files...").pack(pady=10)
    progress_bar = ttk.Progressbar(progress_win, orient='horizontal', length=300, mode='determinate')
    progress_bar.pack(pady=10)
    total_files = len(file_list)
    
    for i, file_path in enumerate(file_list):
        try:
            # 기본값: 파일이 들어있는 폴더 이름을 샘플 이름으로 사용 (기존 txt 방식)
            sample_folder_name = os.path.basename(os.path.dirname(file_path))
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_name)[1].lower()
            
            # 파일 읽기
            encodings_to_try = ['utf-8-sig', 'cp949', 'latin-1', 'utf-8']
            file_content = None
            for enc in encodings_to_try:
                try:
                    with open(file_path, 'r', encoding=enc) as f:
                        file_content = f.readlines()
                    break 
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    print(f"Error reading '{file_name}' with encoding {enc}: {e}")
                    continue

            if file_content is None:
                continue

            jv_data = None
            parameters = {}
            scan_direction = 'Unknown'

            # --- [CSV 파일 처리 로직 변경] ---
            if file_ext == '.csv':
                jv_data, parameters = parse_new_csv_format(file_content)
                if jv_data is None: continue
                
                # [수정됨] 파일명 파싱 로직 제거
                # 이제 파일명이 아닌, 상위 폴더 이름(sample_folder_name)을 그대로 샘플 이름으로 사용합니다.
                # sample_folder_name = os.path.basename(os.path.dirname(file_path)) <- 이미 위에서 설정됨


                if 'Scan Direction' in parameters:
                    scan_direction = parameters['Scan Direction']
                else:
                    scan_direction = detect_scan_direction(file_path, jv_data)

            # --- [TXT 파일 처리 (기존 유지)] ---
            elif file_ext == '.txt':
                data_lines, header_line = [], ""
                header_line = file_content[0].strip()
                line_iter = iter(file_content[1:])
                for line in line_iter:
                    line = line.strip()
                    if not line: continue
                    if line == 'end': break
                    data_lines.append(line)
                
                for line in line_iter: 
                    if ':' in line:
                        parts = line.split(':', 1)
                        if len(parts) == 2: parameters[parts[0].strip()] = parts[1].strip()

                clean_header = '\t'.join(header_line.split('\t')[:3])
                full_data_string = clean_header + '\n' + '\n'.join(data_lines)
                jv_data = pd.read_csv(io.StringIO(full_data_string), sep='\t')
                
                scan_direction = detect_scan_direction(file_path, jv_data)
            
            # --- 공통 데이터 취합 ---
            if jv_data is not None:
                jv_data_cache[file_path] = jv_data

                Rs, Rsh = calculate_resistances(jv_data)
                voc_v = parse_value(parameters.get('Voc (V)'))
                jsc_a_cm2 = parse_value(parameters.get('Jsc(A/cm2)'))
                ff_percent = parse_value(parameters.get('Fill factor (%)'))
                pce_percent = parse_value(parameters.get('Efficiency (%)'))
                
                jsc_ma_cm2 = None
                if isinstance(jsc_a_cm2, (int, float)): 
                    jsc_ma_cm2 = jsc_a_cm2 * 1000.0
                
                all_data.append({
                    'Sample': sample_folder_name, # 여기서 추출한 이름이 그룹명이 됩니다.
                    'File': file_name, 
                    'FullPath': file_path,
                    'Voc (V)': voc_v, 
                    'Jsc (mA/cm2)': jsc_ma_cm2, 
                    'FF (%)': ff_percent, 
                    'PCE (%)': pce_percent,
                    'Rs (Ω·cm²)': Rs, 
                    'Rsh (Ω·cm²)': Rsh, 
                    'Scan Direction': scan_direction
                })

        except Exception as e:
            print(f"Error processing '{file_path}': {e}")
        finally:
            progress_bar['value'] = (i + 1) / total_files * 100
            progress_win.update_idletasks()

    progress_win.destroy()

    if not all_data:
         messagebox.showerror("Error", "No valid data (txt or csv) could be processed.")
         return

    original_all_devices_df = pd.DataFrame(all_data)

    cols_to_numeric = ['Voc (V)', 'Jsc (mA/cm2)', 'FF (%)', 'PCE (%)', 'Rs (Ω·cm²)', 'Rsh (Ω·cm²)']
    for col in cols_to_numeric:
        if col in original_all_devices_df.columns:
            original_all_devices_df[col] = pd.to_numeric(original_all_devices_df[col], errors='coerce')

    current_display_df = original_all_devices_df.copy()
    refresh_all_views(current_display_df)


def merge_selected_folders():
    global current_display_df
    if current_display_df.empty: messagebox.showwarning("No Data", "Please load a folder first."); return

    selected_items = file_tree.selection()
    if not selected_items: messagebox.showwarning("Selection Needed", "Please select at least two sample folders from the list to merge."); return

    folders_to_merge = []
    for item_id in selected_items:
        if not file_tree.parent(item_id): folders_to_merge.append(file_tree.item(item_id, 'text'))

    if len(folders_to_merge) < 2: messagebox.showwarning("Selection Needed", "Please select at least two sample folders to merge."); return
    new_sample_name = simpledialog.askstring("Merge Samples", "Enter a new name for the merged sample group:", parent=root)

    if not new_sample_name or not new_sample_name.strip(): return
    if new_sample_name in current_display_df['Sample'].unique():
        if not messagebox.askyesno("Name Exists", f"The sample name '{new_sample_name}' already exists. Merge selected folders into this existing group?"):
            return

    # IMPORTANT: Update the 'Sample' column in the main display DataFrame
    current_display_df.loc[current_display_df['Sample'].isin(folders_to_merge), 'Sample'] = new_sample_name
    refresh_all_views(current_display_df)


# --- NEW: Auto Merge Folders Function [MODIFIED] ---
def auto_merge_folders():
    global current_display_df
    if current_display_df.empty:
        messagebox.showwarning("No Data", "Please load a folder first.")
        return

    # [수정] 사용자에게 팝업창을 띄워 구분자 입력받기
    delimiter = simpledialog.askstring(
        "Auto Merge by Delimiter",
        "샘플 이름과 숫자를 구분하는 문자를 입력하세요.\n"
        "(예: _  또는  -  또는  (공백) )",
        parent=root
    )

    # 사용자가 'Cancel'을 누른 경우 (None)
    if delimiter is None:
        return

    # [수정] 입력된 구분자에 따라 동적 정규식(Regex) 생성
    separator_pattern = ""
    if delimiter == "" or delimiter.isspace():
        # 사용자가 공백만 입력했거나, 그냥 엔터(빈 문자열)를 친 경우
        # 'A 1', 'B 2' 등을 찾음
        separator_pattern = r"(\s+)" # 1개 이상의 공백
    else:
        # 사용자가 입력한 특수문자(예: '_', '-')를 정규식에서 안전하게 처리 (re.escape)
        # 앞뒤로 0개 이상의 공백을 허용 (예: 'A_1', 'A - 1', 'A-1' 모두 인식)
        escaped_delimiter = re.escape(delimiter)
        separator_pattern = rf"(\s*{escaped_delimiter}\s*)"

    # 최종 정규 표현식: (그룹 이름)(사용자 정의 구분자)(숫자)
    # 예: ^(.*?)(\s*_\s*)(\d+)$
    pattern = re.compile(rf"^(.*?)({separator_pattern})(\d+)$")


    # Use the currently displayed samples for merging logic
    sample_names = current_display_df['Sample'].unique()
    groups_to_merge = defaultdict(list) # 그룹 이름별로 원본 샘플 이름 리스트 저장

    for name in sample_names:
        match = pattern.match(name)
        if match:
            # group(1)이 '그룹 이름', group(3)이 '사용자 정의 구분자', group(4)가 '숫자'가 됩니다.
            # (전체 패턴, 그룹1, (구분자 그룹), 구분자, 그룹4)
            base_name = match.group(1).strip() # 그룹 이름 (앞뒤 공백 제거)
            if base_name: # 그룹 이름이 비어있지 않은 경우만
                groups_to_merge[base_name].append(name)

    merges_performed = 0
    merge_details = []

    if not groups_to_merge:
        # [수정] 메시지에 사용자 입력값 반영
        user_delimiter_msg = "공백" if (delimiter == "" or delimiter.isspace()) else f"'{delimiter}'"
        messagebox.showinfo("Auto Merge", 
                            f"{user_delimiter_msg} 구분자를 사용한 'Name-Number' 패턴의\n"
                            f"샘플 그룹을 찾을 수 없습니다.",
                            parent=root)
        return

    # Create a temporary DataFrame to modify 'Sample' names
    temp_df = current_display_df.copy()

    for base_name, original_names in groups_to_merge.items():
        if len(original_names) > 1: # 병합할 항목이 2개 이상인 경우
            # 데이터프레임에서 해당 샘플들의 이름을 base_name으로 변경
            temp_df.loc[temp_df['Sample'].isin(original_names), 'Sample'] = base_name
            merges_performed += 1
            merge_details.append(f"  - '{base_name}': {', '.join(original_names)}")

    if merges_performed > 0:
         # Update the main DataFrame
        current_display_df = temp_df
        refresh_all_views(current_display_df)
        message = f"Auto-merge complete. {merges_performed} group(s) merged:\n\n"
        message += "\n".join(merge_details)
        messagebox.showinfo("Auto Merge Complete", message)
    else:
        # [수정] 메시지에 사용자 입력값 반영
        user_delimiter_msg = "공백" if (delimiter == "" or delimiter.isspace()) else f"'{delimiter}'"
        messagebox.showinfo("Auto Merge", 
                            f"{user_delimiter_msg} 구분자를 사용한 샘플이 있으나,\n"
                            f"병합할 그룹(2개 이상)이 없습니다.",
                            parent=root)
# --- END NEW ---


# --- [추가] 각 폴더별 스캔 방향에 따른 Best PCE만 남기는 필터 함수 ---
def filter_best_pce_per_sample_direction():
    """
    각 샘플(Sample) 폴더 내에서, 'Device Number'(파일명의 첫 번째 _ 앞 부분)와
    스캔 방향(Forward/Reverse) 별로 PCE가 가장 높은 데이터 1개씩만 남깁니다.
    예: '1_fwd', '1_rev', '2_fwd' ...
    """
    global current_display_df
    if current_display_df.empty:
        messagebox.showwarning("No Data", "Please load data first.")
        return

    try:
        # 1. Device 구분자 추출 (파일명의 첫 번째 '_' 앞 부분)
        # 예: '1_01_fwd.csv' -> '1', 'control_02.csv' -> 'control'
        def get_device_prefix(filename):
            return filename.split('_')[0].strip()

        # 임시 컬럼 생성
        current_display_df['DeviceID'] = current_display_df['File'].apply(get_device_prefix)

        # 2. Sample(폴더), DeviceID(디바이스), Scan Direction 별로 그룹화 -> PCE 최대값 인덱스 찾기
        best_indices = current_display_df.groupby(['Sample', 'DeviceID', 'Scan Direction'])['PCE (%)'].idxmax()
        
        # 3. 해당 인덱스의 행들만 추출
        filtered_df = current_display_df.loc[best_indices].copy()
        
        # 임시 컬럼 삭제
        filtered_df.drop(columns=['DeviceID'], inplace=True)
        current_display_df.drop(columns=['DeviceID'], inplace=True) # 원본 df에서도 삭제

        # 4. 화면 갱신
        refresh_all_views(filtered_df)
        
        # 결과 메시지
        messagebox.showinfo("Best PCE Filter Applied", 
                            f"각 폴더 내 Device 및 스캔 방향별 최고 PCE 데이터만 남겼습니다.\n\n"
                            f"기준: 파일명의 첫 번째 '_' 앞부분을 같은 Device로 인식\n"
                            f"변경 전: {len(current_display_df)} 개\n"
                            f"변경 후: {len(filtered_df)} 개")

    except Exception as e:
        # 혹시 모를 에러 발생 시 임시 컬럼 정리
        if 'DeviceID' in current_display_df.columns:
            current_display_df.drop(columns=['DeviceID'], inplace=True)
        messagebox.showerror("Error", f"Failed to filter best PCE data:\n{e}")
# --- [추가 끝] ---

# --- [신규] 변수 '레시피' 저장/불러오기 함수 ---

def save_variable_recipe():
    """
    현재 선택된 샘플의 7가지 변수 값과, 
    그 값에 연결된 'Process Details'까지 함께 .jvr 파일로 저장합니다.
    """
    global process_details
    
    # 1. 현재 선택된 샘플이 있는지 확인
    sample_name = current_sample_label.get()
    if sample_name == "[No sample selected]":
        messagebox.showwarning("No Sample Selected", 
                             "먼저 변수 '레시피'의 원본이 될 샘플을 좌측 리스트에서 선택하세요.", 
                             parent=root)
        return

    # 2. 저장 파일명 정하기 (신규 확장자 .jvr 사용)
    filepath = filedialog.asksaveasfilename(
        title="Save Variable Recipe",
        defaultextension=".jvr",
        filetypes=[("JV Recipe Files", "*.jvr"), ("All files", "*.*")],
        parent=root
    )
    if not filepath:
        return

    try:
        recipe_data = {
            'variables': {}, # 7가지 변수 값 (문자열)
            'details': {}    # 변수 값에 연결된 공정 상세정보
        }

        # 3. 현재 GUI의 7가지 변수 값(문자열)과 공정 상세정보를 수집
        for col in variable_columns:
            # 3-1. GUI에서 변수 값(문자열) 가져오기
            value_str = var_entry_widgets[col][0].get()
            recipe_data['variables'][col] = value_str
            
            # 3-2. 이 변수 값(예: "NiOx + BCP")을 개별 재료로 분리
            materials = [m.strip() for m in value_str.split(' + ') if m.strip()]
            
            if materials:
                # 3-3. 원본 샘플의 공정 상세정보 딕셔너리 가져오기
                source_details_for_col = process_details.get(sample_name, {}).get(col, {})
                recipe_data['details'][col] = {}
                
                # 3-4. 해당 재료의 상세정보만 recipe_data에 복사
                for material in materials:
                    material_data = source_details_for_col.get(material)
                    if material_data:
                        recipe_data['details'][col][material] = material_data.copy()

        # 4. 수집한 데이터를 .jvr 파일로 저장 (pickle 사용)
        with open(filepath, 'wb') as f:
            pickle.dump(recipe_data, f)
            
        messagebox.showinfo("Success", f"Variable recipe saved to:\n{filepath}", parent=root)

    except Exception as e:
        messagebox.showerror("Save Error", f"Failed to save recipe:\n{e}", parent=root)


def load_variable_recipe():
    """
    .jvr 파일에서 '레시피'를 불러와, 
    현재 선택된 샘플의 GUI 창 값을 덮어쓰고, 공정 상세정보를 임시 적용합니다.
    """
    global process_details
    
    # 1. 레시피를 적용할 대상 샘플이 있는지 확인
    sample_name = current_sample_label.get()
    if sample_name == "[No sample selected]":
        messagebox.showwarning("No Sample Selected", 
                             "먼저 '레시피'를 적용할 샘플을 좌측 리스트에서 선택하세요.", 
                             parent=root)
        return

    # 2. 불러올 .jvr 파일 선택
    filepath = filedialog.askopenfilename(
        title="Load Variable Recipe",
        filetypes=[("JV Recipe Files", "*.jvr"), ("All files", "*.*")],
        parent=root
    )
    if not filepath:
        return

    try:
        # 3. .jvr 파일에서 레시피 데이터 로드
        with open(filepath, 'rb') as f:
            recipe_data = pickle.load(f)

        if not isinstance(recipe_data, dict) or 'variables' not in recipe_data:
             messagebox.showerror("Load Error", "This is not a valid recipe file.", parent=root)
             return
        
        loaded_vars = recipe_data.get('variables', {})
        loaded_details = recipe_data.get('details', {})

        # --- [문제 2 해결] ---
        # 4. (중요) 불러온 변수 값(문자열)을 GUI 엔트리 창에 즉시 덮어쓰기
        for col in variable_columns:
            value_str = loaded_vars.get(col, "")
            var_entry_widgets[col][0].set(value_str)
        
        # 5. (중요) 불러온 공정 상세정보를 현재 선택된 샘플에 덮어쓰기 (메모리에 임시 적용)
        if sample_name not in process_details:
            process_details[sample_name] = {}
            
        for col, materials_dict in loaded_details.items():
            if col not in process_details[sample_name]:
                process_details[sample_name][col] = {}
            
            # 5-1. 해당 'col'의 상세정보를 통째로 덮어씁니다.
            # (이렇게 하면 레시피에 없는 재료의 상세정보는 사라집니다)
            process_details[sample_name][col] = materials_dict.copy()

        messagebox.showinfo("Recipe Loaded", 
                            f"Recipe loaded into '{sample_name}'.\n\n"
                            f"The variable fields and their process details have been updated.\n"
                            f"**Please press 'Save Variables' to apply this recipe permanently.**", 
                            parent=root)

    except Exception as e:
        messagebox.showerror("Load Error", f"Failed to load recipe:\n{e}", parent=root)

# --- Save/Load State Functions ---
def save_state():
    # [수정] operator_name, device_structure 전역 변수 추가
    global current_display_df, experimental_variables, current_root_folder, \
           operator_name, device_structure
           
    if not current_root_folder: messagebox.showwarning("No Data", "Please load a folder first before saving the state."); return
    filepath = filedialog.asksaveasfilename(title="Save Analysis State", defaultextension=".pkl",
                                            filetypes=[("Analysis State Files", "*.pkl"), ("All files", "*.*")])
    if not filepath: return

    # --- 마커/스타일 정보도 함께 저장 ---
    jv_plot_info = {
        item_id: {
            'path': data['path'], 'label': data['label'], 'color': data['color'],
            'marker': data.get('marker', 'o'), 'linestyle': data.get('linestyle', '-') # 호환성을 위해 get 사용
        } for item_id, data in plotted_jv_items.items()
    }

    dist_plot_info = {item_id: {'label': data['label'], 'color': data['color']} for item_id, data in plotted_dist_items.items()}

    # Save the merge state based on the current display df
    merge_info = pd.Series(current_display_df.Sample.values, index=current_display_df.FullPath).to_dict()


    var_check_states = {col: var_checkbox_vars[col].get() for col in variable_columns if col in var_checkbox_vars} # Check if key exists

    state = {
        'root_folder': current_root_folder,
        'operator_name': operator_name, # [신규] 사용자 이름 저장
        'device_structure': device_structure, # [신규] 소자 구조 저장
        'jv_plots': jv_plot_info,
        'dist_plots': dist_plot_info,
        'filters': { "voc_min": filter_voc_min.get(), "voc_max": filter_voc_max.get(),
                     "jsc_min": filter_jsc_min.get(), "jsc_max": filter_jsc_max.get(),
                     "ff_min": filter_ff_min.get(), "ff_max": filter_ff_max.get(),
                     "pce_min": filter_pce_min.get(), "pce_max": filter_pce_max.get() },
        'merge_info': merge_info, # Save current merge state
        'scan_filters': { 'jv': scan_filter_jv_var.get(), 'dist': scan_filter_dist_var.get() },
        'experimental_variables': experimental_variables,
        'var_check_states': var_check_states,
        'process_details': process_details
    }

    try:
        with open(filepath, 'wb') as f:
            pickle.dump(state, f)
        messagebox.showinfo("Success", f"Analysis state saved to:\n{filepath}")
    except Exception as e:
        messagebox.showerror("Save Error", f"Failed to save state:\n{e}")

def load_state():
    # [수정] operator_name, device_structure 전역 변수 추가
    global current_display_df, experimental_variables, original_all_devices_df, jv_data_cache, \
           process_details, operator_name, device_structure
           
    filepath = filedialog.askopenfilename(title="Load Analysis State", filetypes=[("Analysis State Files", "*.pkl"), ("All files", "*.*")])
    if not filepath: return

    try:
        with open(filepath, 'rb') as f:
            state = pickle.load(f)

        # --- 수정된 부분 (1) ---
        # 불러온 변수를 전역 변수가 아닌 로컬 임시 변수에 저장합니다.
        # (load_and_process_folder가 전역 변수를 초기화하기 때문)
        loaded_exp_vars = state.get('experimental_variables', {})
        loaded_proc_details = state.get('process_details', {})
        
        # --- [신규] 사용자 이름과 소자 구조를 먼저 불러옵니다. ---
        # load_and_process_folder가 팝업을 띄우지 않도록 하기 위함
        operator_name = state.get('operator_name', '')
        device_structure = state.get('device_structure', '')
        # --- [신규 끝] ---


        # 1. Load data from the saved root folder first
        # [수정] 이 함수는 is_fresh_load=False로 실행되어 팝업을 띄우지 않습니다.
        load_and_process_folder(state['root_folder']) # This resets current_display_df to original

        # --- 수정된 부분 (2) ---
        # load_and_process_folder에 의해 초기화된 전역 변수들을
        # 아까 임시 저장해둔 로컬 변수 값으로 다시 복원합니다.
        experimental_variables = loaded_exp_vars
        process_details = loaded_proc_details
        # --- 수정 끝 ---

        # 2. Apply the saved merge state AFTER loading
        merge_info = state.get('merge_info', {})
        if merge_info:
            # Apply merge map to the newly loaded current_display_df
            current_display_df['Sample'] = current_display_df['FullPath'].map(merge_info).fillna(current_display_df['Sample'])

        # 3. Apply filters
        filters = state.get('filters', {})
        filter_voc_min.set(filters.get("voc_min", "")); filter_voc_max.set(filters.get("voc_max", ""))
        filter_jsc_min.set(filters.get("jsc_min", "")); filter_jsc_max.set(filters.get("jsc_max", ""))
        filter_ff_min.set(filters.get("ff_min", "")); filter_ff_max.set(filters.get("ff_max", ""))
        filter_pce_min.set(filters.get("pce_min", "")); filter_pce_max.set(filters.get("pce_max", ""))
        apply_filter() # Apply filter values to the merged data


        # 4. Set Scan filters
        scan_filters = state.get('scan_filters', {})
        scan_filter_jv_var.set(scan_filters.get('jv', 'All'))
        scan_filter_dist_var.set(scan_filters.get('dist', 'All'))

        # 5. Set Variable Check states
        var_check_states = state.get('var_check_states', {})
        for col, is_checked in var_check_states.items():
            if col in var_checkbox_vars:
                var_checkbox_vars[col].set(is_checked)


        # 6. Refresh UI completely with merged and filtered data
        # apply_filter() already calls refresh_all_views
        # refresh_all_views(current_display_df)


        # 7. Restore Plots
        clear_jv_plot()
        for item_id, plot_info in state.get('jv_plots', {}).items():
             # Check if the file path exists in the current (potentially filtered) data
            if plot_info['path'] in current_display_df['FullPath'].values:
                df = jv_data_cache.get(plot_info['path'])
                if df is not None:
                    color_img = create_color_image(plot_info['color'])
                    new_id = plotted_list_tree.insert('', 'end', text=plot_info['label'], image=color_img)
                    plotted_jv_items[new_id] = {
                        'path': plot_info['path'], 'label': plot_info['label'],
                        'color': plot_info['color'], 'df': df, 'image': color_img,
                        'marker': plot_info.get('marker', 'o'),
                        'linestyle': plot_info.get('linestyle', '-'),
                        # Restore scan direction info if needed for labels, etc.
                        'scan_dir': current_display_df.loc[current_display_df['FullPath'] == plot_info['path'], 'Scan Direction'].iloc[0]
                    }
            else:
                 print(f"Warning: J-V plot item for '{plot_info['path']}' not found in current data after loading state. Skipping.")
        redraw_jv_graphs()

        clear_dist_plot()
        for item_id, plot_info in state.get('dist_plots', {}).items():
            # Check if the folder name (label) still exists after merge/filter
            if plot_info['label'] in current_display_df['Sample'].unique():
                color_img = create_color_image(plot_info['color'])
                new_id = dist_list_tree.insert('', 'end', text=plot_info['label'], image=color_img)
                plotted_dist_items[new_id] = {'label': plot_info['label'], 'color': plot_info['color'], 'image': color_img}
            else:
                print(f"Warning: Distribution plot item '{plot_info['label']}' not found in current data after loading state. Skipping.")
        redraw_dist_plot()


    except FileNotFoundError:
         messagebox.showerror("Load Error", f"Failed to load state: Original data folder not found at\n{state.get('root_folder', 'Unknown Path')}")
    except Exception as e:
        messagebox.showerror("Load Error", f"Failed to load state:\n{e}")


# --- 유틸리티 함수 ---
def copy_treeview_to_clipboard(treeview_widget):
    headers = [treeview_widget.heading(col)['text'] for col in treeview_widget['columns']]
    header_line = '\t'.join(headers)
    data_lines = []
    for item_id in treeview_widget.get_children():
        values = treeview_widget.item(item_id)['values']
        data_lines.append('\t'.join(map(str, values)))
    full_text = header_line + '\n' + '\n'.join(data_lines)
    root.clipboard_clear()
    root.clipboard_append(full_text)
    messagebox.showinfo("Copied", "Table data has been copied to the clipboard.\n\nFor PowerPoint, it's recommended to paste into Excel first, then copy the table from Excel to PowerPoint.")

def on_closing():
    if messagebox.askokcancel("Quit", "Do you want to quit?"):
        plt.close('all') # Close matplotlib figures explicitly
        root.quit()
        root.destroy()

# --- [수정] Perovskite Builder Window (Additive 추가) ---
class PerovskiteBuilderWindow(tk.Toplevel):
    def __init__(self, parent, sample_name):
        super().__init__(parent)
        self.title(f"Perovskite Builder for {sample_name}")
        self.geometry("500x500") # 높이 증가
        self.transient(parent) # Keep on top
        self.grab_set() # Modal

        self.sample_name = sample_name
        self.result_var = var_entry_widgets["Perovskite"][0]

        # Load existing data if possible
        self.var_data = experimental_variables.get(self.sample_name, {})

        main_frame = ttk.Frame(self, padding=10)
        main_frame.pack(fill='both', expand=True)

        self.site_vars = {}
        self.a_ions = ['Cs', 'FA', 'MA']; self.b_ions = ['Pb', 'Sn']; self.x_ions = ['I', 'Br', 'Cl']

        a_frame = ttk.LabelFrame(main_frame, text="A-Site")
        a_frame.pack(fill='x', expand=True, pady=5, padx=5)
        self.site_vars['A'] = self._create_site_entries(a_frame, self.a_ions)

        b_frame = ttk.LabelFrame(main_frame, text="B-Site")
        b_frame.pack(fill='x', expand=True, pady=5, padx=5)
        self.site_vars['B'] = self._create_site_entries(b_frame, self.b_ions)

        x_frame = ttk.LabelFrame(main_frame, text="X-Site")
        x_frame.pack(fill='x', expand=True, pady=5, padx=5)
        self.site_vars['X'] = self._create_site_entries(x_frame, self.x_ions)

        # --- [추가] Additive Frame ---
        additive_frame = ttk.LabelFrame(main_frame, text="Additive")
        additive_frame.pack(fill='x', expand=True, pady=5, padx=5)
        
        add_conc_frame = ttk.Frame(additive_frame)
        add_conc_frame.pack(fill='x', pady=3, padx=5)
        ttk.Label(add_conc_frame, text="Name:", width=8).pack(side=tk.LEFT)
        self.additive_name_var = tk.StringVar()
        ttk.Entry(add_conc_frame, textvariable=self.additive_name_var, width=15).pack(side=tk.LEFT, padx=5)
        
        ttk.Label(add_conc_frame, text="Conc:", width=5).pack(side=tk.LEFT)
        self.additive_conc_var = tk.StringVar()
        ttk.Entry(add_conc_frame, textvariable=self.additive_conc_var, width=8).pack(side=tk.LEFT, padx=5)
        self.additive_unit_var = tk.StringVar(value="mg/mL") # Default unit
        ttk.Combobox(add_conc_frame, textvariable=self.additive_unit_var, values=["mg/mL", "M", "mol%", "vol%", "wt%"], width=6).pack(side=tk.LEFT)
        # --- [추가 끝] ---

        # --- Bandgap Entry ---
        bandgap_frame = ttk.Frame(main_frame)
        bandgap_frame.pack(fill='x', padx=10, pady=5)
        
        ttk.Label(bandgap_frame, text="Bandgap ($E_g$):", width=15).pack(side=tk.LEFT)
        self.bandgap_var = tk.StringVar() 
        bandgap_entry = ttk.Entry(bandgap_frame, textvariable=self.bandgap_var, width=10)
        bandgap_entry.pack(side=tk.LEFT, padx=5)
        ttk.Label(bandgap_frame, text="eV").pack(side=tk.LEFT)
        # --- End of Bandgap Entry ---

        ttk.Label(main_frame, text="Generated Formula:").pack(pady=(10,0))
        self.formula_var = tk.StringVar()
        ttk.Entry(main_frame, textvariable=self.formula_var, state='readonly').pack(fill='x', expand=True, padx=5)

        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(pady=10)
        ttk.Button(btn_frame, text="Generate & Apply", command=self.generate_and_apply).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_frame, text="Cancel", command=self.destroy).pack(side=tk.LEFT, padx=5)

        self._load_initial_values()

    def _create_site_entries(self, parent_frame, ions):
        vars = {}
        for ion in ions:
            frame = ttk.Frame(parent_frame)
            chk_var = tk.BooleanVar()
            chk = ttk.Checkbutton(frame, text=ion, variable=chk_var)
            val_var = tk.StringVar(value="0.0")
            entry = ttk.Entry(frame, textvariable=val_var, width=5)
            chk.pack(side=tk.LEFT, padx=(10, 2))
            entry.pack(side=tk.LEFT, padx=(0, 10))
            frame.pack(side=tk.LEFT)
            vars[ion] = (chk_var, val_var)
        return vars

    def _load_initial_values(self):
        formula = experimental_variables.get(self.sample_name, {}).get("Perovskite", "")
        self.formula_var.set(formula)
        
        # --- [수정] Parse Additive and Bandgap ---
        # Regex for Additive: + Name (Conc Unit)
        add_match = re.search(r"\+\s*([\w\d\.-]+)\s*\((\S+)\s*([\w\/%]+)\)", formula)
        if add_match:
            self.additive_name_var.set(add_match.group(1))
            self.additive_conc_var.set(add_match.group(2))
            self.additive_unit_var.set(add_match.group(3))
        else:
            # Simpler additive match (just name, no conc/unit)
            # A, B, X 사이트 파싱 전에 실행되어야 함
            base_formula = formula.split(' (Eg=')[0] # Eg 파트 제거
            add_match_simple = re.search(r"\+\s*([\w\d\.-]+)", base_formula)
            if add_match_simple:
                 self.additive_name_var.set(add_match_simple.group(1).strip())
        
        # Regex for Bandgap: (Eg=...eV)
        eg_match = re.search(r"\(Eg=([\d\.]+)\s*eV\)", formula)
        if eg_match:
            self.bandgap_var.set(eg_match.group(1))
        # --- [수정 끝] ---

        # Simplified parsing logic
        for ion_list, site_key in [(self.a_ions, 'A'), (self.b_ions, 'B'), (self.x_ions, 'X')]:
            for ion in ion_list:
                # --- MODIFIED: More robust parsing ---
                # Use regex to find the ion possibly followed by numbers/dots
                pattern = rf"(?<![A-Za-z])({ion})(?:([\d\.]+))?"
                match = re.search(pattern, formula)
                if match:
                    self.site_vars[site_key][ion][0].set(True)
                    stoichiometry = match.group(2)
                    if stoichiometry:
                        try:
                            float_val = float(stoichiometry)
                            self.site_vars[site_key][ion][1].set(f"{float_val:.2f}")
                        except ValueError:
                            self.site_vars[site_key][ion][1].set("1.0") # Default if parsing fails
                    else:
                         self.site_vars[site_key][ion][1].set("1.0") # Default if no number follows
                # --- END MODIFIED ---


    def generate_and_apply(self):
        def format_site(ions_vars):
            parts = []
            total_stoich = 0
            for ion, (chk_var, val_var) in ions_vars.items():
                if chk_var.get():
                    try:
                        val = float(val_var.get())
                        if val > 0:
                            parts.append((ion, val))
                            total_stoich += val
                    except ValueError:
                        parts.append((ion, -1)) # Mark as invalid

            # Normalize if total stoichiometry is given and > 0, otherwise just list
            site_str = ""
            if total_stoich > 0 and abs(total_stoich - 1.0) > 1e-6 : # Normalize if not close to 1
                 # Normalize only if multiple ions are selected
                 if len(parts) > 1:
                    site_str = "".join([f"{ion}{val/total_stoich:.2f}".rstrip('0').rstrip('.') for ion, val in parts if val != -1])
                 else: # If only one ion, don't normalize to 1 unless it was 1 originally
                     ion, val = parts[0]
                     if abs(val - 1.0) < 1e-6: site_str = ion
                     else: site_str = f"{ion}{val:.2f}".rstrip('0').rstrip('.')

            else: # If total is close to 1 or cannot normalize
                site_str = "".join([ion if abs(val-1.0) < 1e-6 else (f"{ion}{val:.2f}".rstrip('0').rstrip('.') if val > 0 else f"{ion}?") for ion, val in parts])


            # Wrap with parentheses if multiple components or single component needs explicit stoichiometry
            if len(parts) > 1:
                return f"({site_str})"
            elif len(parts) == 1 and "?" not in site_str and not site_str.isalpha(): # Single ion with number
                return f"({site_str})"
            else:
                 return site_str # Single ion like 'FA' or 'Pb'

        a_str = format_site(self.site_vars['A'])
        b_str = format_site(self.site_vars['B'])
        x_str = format_site(self.site_vars['X'])

        if not a_str: a_str = "FA" # Default
        if not b_str: b_str = "Pb" # Default
        if not x_str: x_str = "I"  # Default

        # Construct formula, handle potential empty B site
        if b_str:
            full_formula = f"{a_str}{b_str}{x_str}3"
        else: # Handle case B site might be empty (though default prevents this)
             full_formula = f"{a_str}(?){x_str}3"

        # --- [수정] Append Additive and Bandgap ---
        additive_name = self.additive_name_var.get().strip()
        additive_conc = self.additive_conc_var.get().strip()
        additive_unit = self.additive_unit_var.get().strip()
        
        if additive_name and additive_conc:
            full_formula += f" + {additive_name} ({additive_conc} {additive_unit})"
        elif additive_name:
            full_formula += f" + {additive_name}"

        bandgap_value = self.bandgap_var.get().strip()
        if bandgap_value:
            try:
                # Validate it's a number
                float(bandgap_value)
                full_formula += f" (Eg={bandgap_value}eV)"
            except ValueError:
                pass # Don't append if not a valid number
        # --- [수정 끝] ---

        self.formula_var.set(full_formula)
        self.result_var.set(full_formula) # Update main window's entry
        self.destroy()

# --- [수정] 다중 선택 팝업창 클래스 (Custom Material 추가) ---
class EnhancedVariableSelectorWindow(tk.Toplevel):
    """[수정] 다중 선택 + 공정 조건 입력 + '커스텀 입력'이 가능한 변수 선택 창"""
    
    def __init__(self, parent, title, variable_name, options, result_var):
        super().__init__(parent)
        self.title(title)
        self.geometry("700x600") # 높이 증가
        self.transient(parent)
        self.grab_set()

        self.variable_name = variable_name
        self.options = options
        self.result_var = result_var
        self.vars = {}
        self.detail_buttons = {}
        
        self.current_sample = current_sample_label.get()
        current_values = [v.strip() for v in self.result_var.get().split(' + ') if v.strip()]

        main_frame = ttk.Frame(self, padding=10)
        main_frame.pack(fill="both", expand=True)
        
        # --- Presets Frame ---
        preset_frame = ttk.LabelFrame(main_frame, text="Presets", padding=10)
        preset_frame.pack(fill="both", expand=True)

        canvas = tk.Canvas(preset_frame)
        scrollbar = ttk.Scrollbar(preset_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = ttk.Frame(canvas)

        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )

        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)

        for i, option in enumerate(self.options):
            option_frame = ttk.Frame(scrollable_frame)
            option_frame.pack(fill='x', padx=5, pady=3)
            
            var = tk.BooleanVar()
            chk = ttk.Checkbutton(option_frame, text=option, variable=var, width=15)
            
            if option in current_values:
                var.set(True)
            
            chk.pack(side=tk.LEFT, padx=5)
            self.vars[option] = var
            
            detail_btn = ttk.Button(
                option_frame, 
                text="Details...",
                command=lambda opt=option: self.open_detail_window(opt),
                width=10
            )
            detail_btn.pack(side=tk.LEFT, padx=5)
            self.detail_buttons[option] = detail_btn

        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # --- [추가] Custom Material Input Frame ---
        custom_frame = ttk.LabelFrame(main_frame, text="Custom Material", padding=10)
        custom_frame.pack(fill='x', padx=5, pady=(10, 0))

        custom_input_frame = ttk.Frame(custom_frame)
        custom_input_frame.pack(fill='x')
        
        ttk.Label(custom_input_frame, text="Name:", width=8).pack(side=tk.LEFT)
        
        self.custom_var = tk.StringVar()
        custom_entry = ttk.Entry(custom_input_frame, textvariable=self.custom_var, width=20)
        custom_entry.pack(side=tk.LEFT, padx=5)
        
        # 프리셋에 없는 커스텀 값 로드 시도
        preset_options = set(self.options)
        custom_val_found = ""
        for val in current_values:
            if val not in preset_options:
                custom_val_found = val # 프리셋에 없는 첫 번째 값을 커스텀 값으로 간주
                break
        self.custom_var.set(custom_val_found)

        custom_detail_btn = ttk.Button(
            custom_input_frame,
            text="Details...",
            command=self.open_custom_detail_window, # 새 함수 연결
            width=10
        )
        custom_detail_btn.pack(side=tk.LEFT, padx=5)
        
        ttk.Label(
            custom_frame, 
            text="프리셋에 없는 물질을 직접 입력하고 'Details...' 버튼으로 공정 조건을 추가하세요."
        ).pack(fill='x', pady=(5,0))
        # --- [추가 끝] ---

        # --- Bottom Button Frame ---
        btn_frame = ttk.Frame(self, padding=10)
        btn_frame.pack(side=tk.BOTTOM, fill='x')

        ttk.Button(btn_frame, text="Apply", command=self.apply).pack(side=tk.RIGHT, padx=5)
        ttk.Button(btn_frame, text="Cancel", command=self.destroy).pack(side=tk.RIGHT)

    def open_detail_window(self, material):
        """(기존) 프리셋 재료에 대한 상세 공정 조건 입력 창"""
        ProcessDetailWindow(self, self.current_sample, self.variable_name, material)

    # --- [추가] 새 함수 추가 ---
    def open_custom_detail_window(self):
        """'커스텀 재료'에 대한 상세 공정 조건 입력 창"""
        custom_material_name = self.custom_var.get().strip()
        if not custom_material_name:
            messagebox.showwarning(
                "No Name", 
                "먼저 커스텀 물질의 이름을 입력하세요.",
                parent=self
            )
            return
        
        # 기존 ProcessDetailWindow 재사용
        ProcessDetailWindow(
            self, 
            self.current_sample, 
            self.variable_name, 
            custom_material_name
        )
    # --- [추가 끝] ---

    # --- [수정] apply 함수 수정 ---
    def apply(self):
        """선택된 항목들 + 커스텀 항목을 " + "로 연결하여 결과 변수에 저장"""
        # 1. Get selected presets
        selected_items = [option for option, var in self.vars.items() if var.get()]
        
        # 2. Get custom material
        custom_material = self.custom_var.get().strip()
        
        # 3. Add custom material if it exists and is not already in the preset list
        if custom_material and custom_material not in selected_items:
            selected_items.append(custom_material)
        
        # 4. Apply
        self.result_var.set(" + ".join(selected_items))
        self.destroy()
    # --- [수정 끝] ---


# --- [수정] 공정 상세정보 입력 창 (2단 레이아웃 적용) ---
class ProcessDetailWindow(tk.Toplevel):
    """공정 조건 상세 입력 창"""
    
    def __init__(self, parent, sample_name, variable_name, material):
        super().__init__(parent)
        self.title(f"{variable_name} - {material} Process Details")
        
        # [수정] 2단 레이아웃을 위해 창 기본 너비 증가
        self.geometry("900x750") 
        
        self.transient(parent)
        self.grab_set()
        
        self.sample_name = sample_name
        self.variable_name = variable_name # TCO, HTL, Contact...
        self.material = material
        
        if sample_name not in process_details:
            process_details[sample_name] = {}
        if variable_name not in process_details[sample_name]:
            process_details[sample_name][variable_name] = {}
        if material not in process_details[sample_name][variable_name]:
            process_details[sample_name][variable_name][material] = {}
            
        self.existing_data = process_details[sample_name][variable_name][material]
        
        main_frame = ttk.Frame(self, padding=15)
        main_frame.pack(fill='both', expand=True)
        
        title_label = ttk.Label(
            main_frame, 
            text=f"{variable_name}: {material}", 
            font=('Helvetica', 12, 'bold')
        )
        title_label.pack(pady=(0, 15))
        
        self.entries = {}
        
        # --- [신규] 2단 레이아웃을 위한 컨테이너 프레임 ---
        # 이 프레임은 'else' 블록에서만 사용됨
        column_container = ttk.Frame(main_frame)
        
        if variable_name == "TCO":
            # --- 1. TCO 전용 UI (1단) ---
            note_frame_parent = main_frame # Notes가 붙을 부모
            
            tco_frame = ttk.LabelFrame(main_frame, text="Substrate Treatment", padding=10)
            tco_frame.pack(fill='x', pady=5, anchor='n')
            # ... (UVO, Ar-Plasma 코드) ...
            uvo_frame = ttk.Frame(tco_frame)
            uvo_frame.pack(fill='x', pady=3)
            ttk.Label(uvo_frame, text="UVO Time:", width=15).pack(side=tk.LEFT)
            uvo_var = tk.StringVar(value=self.existing_data.get('tco_uvo_time', ''))
            uvo_entry = ttk.Entry(uvo_frame, textvariable=uvo_var, width=12)
            uvo_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(uvo_frame, text="min").pack(side=tk.LEFT)
            self.entries['tco_uvo_time'] = uvo_var
            plasma_frame = ttk.Frame(tco_frame)
            plasma_frame.pack(fill='x', pady=3)
            ttk.Label(plasma_frame, text="Ar-Plasma Time:", width=15).pack(side=tk.LEFT)
            plasma_var = tk.StringVar(value=self.existing_data.get('tco_plasma_time', ''))
            plasma_entry = ttk.Entry(plasma_frame, textvariable=plasma_var, width=12)
            plasma_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(plasma_frame, text="min").pack(side=tk.LEFT)
            self.entries['tco_plasma_time'] = plasma_var

        elif variable_name == "Contact":
            # --- 2. Contact 전용 UI (1단) ---
            note_frame_parent = main_frame # Notes가 붙을 부모
            
            contact_frame = ttk.LabelFrame(main_frame, text="Deposition Details", padding=10)
            contact_frame.pack(fill='x', pady=5, anchor='n')
            # ... (Thickness 코드) ...
            thick_frame = ttk.Frame(contact_frame)
            thick_frame.pack(fill='x', pady=3)
            ttk.Label(thick_frame, text="Thickness:", width=15).pack(side=tk.LEFT)
            thick_var = tk.StringVar(value=self.existing_data.get('contact_thickness', ''))
            thick_entry = ttk.Entry(thick_frame, textvariable=thick_var, width=12)
            thick_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(thick_frame, text="nm").pack(side=tk.LEFT)
            self.entries['contact_thickness'] = thick_var

        else:
            # --- 3. 기타 모든 레이어 (HTL, Perovskite, ETL 등)의 표준 UI (2단) ---
            
            # 2단 컨테이너를 main_frame에 pack
            column_container.pack(fill='both', expand=True)
            
            # --- [신규] 왼쪽 / 오른쪽 컬럼 프레임 생성 ---
            left_column = ttk.Frame(column_container)
            left_column.pack(side=tk.LEFT, fill='y', padx=(0, 10), anchor='n')
            
            right_column = ttk.Frame(column_container)
            right_column.pack(side=tk.LEFT, fill='both', expand=True)
            
            note_frame_parent = right_column # Notes가 붙을 부모

            # --- 왼쪽 컬럼 위젯 ---
            
            # (용액 제조 섹션)
            solution_frame = ttk.LabelFrame(left_column, text="Solution Preparation", padding=10)
            solution_frame.pack(fill='x', pady=5)
            # ... (Concentration, Solvents, Volume, Stirring, Aging 코드) ...
            conc_frame = ttk.Frame(solution_frame)
            conc_frame.pack(fill='x', pady=3)
            ttk.Label(conc_frame, text="Concentration:", width=15).pack(side=tk.LEFT)
            init_conc_val = self.existing_data.get('concentration_val', self.existing_data.get('concentration', ''))
            conc_val_var = tk.StringVar(value=init_conc_val)
            conc_entry = ttk.Entry(conc_frame, textvariable=conc_val_var, width=12)
            conc_entry.pack(side=tk.LEFT, padx=5)
            conc_unit_var = tk.StringVar(value=self.existing_data.get('concentration_unit', 'mg/mL'))
            conc_unit_combo = ttk.Combobox(
                conc_frame, textvariable=conc_unit_var, 
                values=["mg/mL", "M"], width=6, state='readonly'
            )
            conc_unit_combo.pack(side=tk.LEFT)
            self.entries['concentration_val'] = conc_val_var
            self.entries['concentration_unit'] = conc_unit_var
            solvent1_frame = ttk.Frame(solution_frame)
            solvent1_frame.pack(fill='x', pady=3)
            ttk.Label(solvent1_frame, text="Solvent 1:", width=15).pack(side=tk.LEFT)
            sol1_name_var = tk.StringVar(value=self.existing_data.get('solvent_1_name', ''))
            sol1_name_entry = ttk.Entry(solvent1_frame, textvariable=sol1_name_var, width=12)
            sol1_name_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(solvent1_frame, text="Ratio 1:").pack(side=tk.LEFT)
            sol1_ratio_var = tk.StringVar(value=self.existing_data.get('solvent_1_ratio', ''))
            sol1_ratio_entry = ttk.Entry(solvent1_frame, textvariable=sol1_ratio_var, width=4)
            sol1_ratio_entry.pack(side=tk.LEFT, padx=5)
            self.entries['solvent_1_name'] = sol1_name_var
            self.entries['solvent_1_ratio'] = sol1_ratio_var
            solvent2_frame = ttk.Frame(solution_frame)
            solvent2_frame.pack(fill='x', pady=3)
            ttk.Label(solvent2_frame, text="Solvent 2:", width=15).pack(side=tk.LEFT)
            sol2_name_var = tk.StringVar(value=self.existing_data.get('solvent_2_name', ''))
            sol2_name_entry = ttk.Entry(solvent2_frame, textvariable=sol2_name_var, width=12)
            sol2_name_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(solvent2_frame, text="Ratio 2:").pack(side=tk.LEFT)
            sol2_ratio_var = tk.StringVar(value=self.existing_data.get('solvent_2_ratio', ''))
            sol2_ratio_entry = ttk.Entry(solvent2_frame, textvariable=sol2_ratio_var, width=4)
            sol2_ratio_entry.pack(side=tk.LEFT, padx=5)
            self.entries['solvent_2_name'] = sol2_name_var
            self.entries['solvent_2_ratio'] = sol2_ratio_var
            solvent3_frame = ttk.Frame(solution_frame)
            solvent3_frame.pack(fill='x', pady=3)
            ttk.Label(solvent3_frame, text="Solvent 3:", width=15).pack(side=tk.LEFT)
            sol2_name_var = tk.StringVar(value=self.existing_data.get('solvent_3_name', ''))
            sol2_name_entry = ttk.Entry(solvent3_frame, textvariable=sol2_name_var, width=12)
            sol2_name_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(solvent3_frame, text="Ratio 3:").pack(side=tk.LEFT)
            sol2_ratio_var = tk.StringVar(value=self.existing_data.get('solvent_3_ratio', ''))
            sol2_ratio_entry = ttk.Entry(solvent3_frame, textvariable=sol2_ratio_var, width=4)
            sol2_ratio_entry.pack(side=tk.LEFT, padx=5)
            self.entries['solvent_3_name'] = sol2_name_var
            self.entries['solvent_3_ratio'] = sol2_ratio_var
            volume_frame = ttk.Frame(solution_frame)
            volume_frame.pack(fill='x', pady=3)
            ttk.Label(volume_frame, text="Drop Volume:", width=15).pack(side=tk.LEFT)
            volume_var = tk.StringVar(value=self.existing_data.get('volume', ''))
            volume_entry = ttk.Entry(volume_frame, textvariable=volume_var, width=12)
            volume_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(volume_frame, text="μL").pack(side=tk.LEFT)
            self.entries['volume'] = volume_var
            stir_temp_frame = ttk.Frame(solution_frame)
            stir_temp_frame.pack(fill='x', pady=3)
            ttk.Label(stir_temp_frame, text="Stirring Temp:", width=15).pack(side=tk.LEFT)
            stir_temp_var = tk.StringVar(value=self.existing_data.get('solution_stir_temp', ''))
            stir_temp_entry = ttk.Entry(stir_temp_frame, textvariable=stir_temp_var, width=12)
            stir_temp_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(stir_temp_frame, text="°C").pack(side=tk.LEFT)
            self.entries['solution_stir_temp'] = stir_temp_var
            stir_time_frame = ttk.Frame(solution_frame)
            stir_time_frame.pack(fill='x', pady=3)
            ttk.Label(stir_time_frame, text="Stirring Time:", width=15).pack(side=tk.LEFT)
            stir_time_val_var = tk.StringVar(value=self.existing_data.get('solution_stir_time_val', ''))
            stir_time_entry = ttk.Entry(stir_time_frame, textvariable=stir_time_val_var, width=12)
            stir_time_entry.pack(side=tk.LEFT, padx=5)
            stir_time_unit_var = tk.StringVar(value=self.existing_data.get('solution_stir_time_unit', 'min'))
            stir_time_combo = ttk.Combobox(
                stir_time_frame, textvariable=stir_time_unit_var, 
                values=["min", "hr"], width=6, state='readonly'
            )
            stir_time_combo.pack(side=tk.LEFT)
            self.entries['solution_stir_time_val'] = stir_time_val_var
            self.entries['solution_stir_time_unit'] = stir_time_unit_var
            aging_frame = ttk.Frame(solution_frame)
            aging_frame.pack(fill='x', pady=3)
            ttk.Label(aging_frame, text="Aging Time:", width=15).pack(side=tk.LEFT)
            aging_var = tk.StringVar(value=self.existing_data.get('solution_aging_time', ''))
            aging_entry = ttk.Entry(aging_frame, textvariable=aging_var, width=12)
            aging_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(aging_frame, text="hr (after prep)").pack(side=tk.LEFT)
            self.entries['solution_aging_time'] = aging_var

            # (증착 방법 - Static/Dynamic)
            dep_method_frame = ttk.LabelFrame(left_column, text="Deposition Method", padding=10)
            dep_method_frame.pack(fill='x', pady=5)
            # ... (Static/Dynamic 코드) ...
            method_frame = ttk.Frame(dep_method_frame)
            method_frame.pack(fill='x', pady=3)
            ttk.Label(method_frame, text="Method:", width=15).pack(side=tk.LEFT)
            self.method_var = tk.StringVar(value=self.existing_data.get('deposition_method', 'Static')) 
            method_combo = ttk.Combobox(
                method_frame, textvariable=self.method_var,
                values=["Static", "Dynamic"], width=10, state='readonly'
            )
            method_combo.pack(side=tk.LEFT, padx=5)
            method_combo.bind('<<ComboboxSelected>>', self._on_deposition_method_selected)
            self.dynamic_drop_time_var = tk.StringVar(value=self.existing_data.get('dynamic_drop_time', ''))
            self.entries['deposition_method'] = self.method_var
            self.entries['dynamic_drop_time'] = self.dynamic_drop_time_var
            self.dep_method_detail_frame = ttk.Frame(dep_method_frame, padding=(10, 5))
            self.dep_method_detail_frame.pack(fill='x', expand=True)

            # (환경 변수)
            env_master_frame = ttk.LabelFrame(left_column, text="Processing Environment", padding=10)
            env_master_frame.pack(fill='x', pady=5)
            # ... (Environment, Temp, O2, H2O, Humidity 코드) ...
            env_frame = ttk.Frame(env_master_frame)
            env_frame.pack(fill='x', pady=3)
            ttk.Label(env_frame, text="Environment:", width=15).pack(side=tk.LEFT)
            self.env_var = tk.StringVar(value=self.existing_data.get('spin_environment', 'Glovebox'))
            self.env_temp_var = tk.StringVar(value=self.existing_data.get('spin_env_temp', ''))
            self.env_humidity_var = tk.StringVar(value=self.existing_data.get('spin_env_humidity', ''))
            self.env_o2_var = tk.StringVar(value=self.existing_data.get('spin_env_o2', ''))
            self.env_h2o_var = tk.StringVar(value=self.existing_data.get('spin_env_h2o', ''))
            self.entries['spin_environment'] = self.env_var
            self.entries['spin_env_temp'] = self.env_temp_var
            self.entries['spin_env_humidity'] = self.env_humidity_var
            self.entries['spin_env_o2'] = self.env_o2_var
            self.entries['spin_env_h2o'] = self.env_h2o_var
            env_combo = ttk.Combobox(
                env_frame, 
                textvariable=self.env_var, 
                values=["Glovebox", "Air"], 
                width=10, 
                state='readonly'
            )
            env_combo.pack(side=tk.LEFT, padx=5)
            env_combo.bind('<<ComboboxSelected>>', self._on_environment_selected)
            self.env_detail_frame = ttk.Frame(env_master_frame, padding=(10, 5))
            self.env_detail_frame.pack(fill='x', expand=True)

            # --- 오른쪽 컬럼 위젯 ---
            
            # (스핀 코팅 1차)
            spin_frame_1 = ttk.LabelFrame(right_column, text="Spin Coating - 1st Step", padding=10)
            spin_frame_1.pack(fill='x', pady=5)
            # ... (RPM 1, Time 1 코드) ...
            rpm_frame_1 = ttk.Frame(spin_frame_1)
            rpm_frame_1.pack(fill='x', pady=3)
            ttk.Label(rpm_frame_1, text="Speed (1st):", width=15).pack(side=tk.LEFT)
            rpm_var_1 = tk.StringVar(value=self.existing_data.get('rpm', ''))
            rpm_entry_1 = ttk.Entry(rpm_frame_1, textvariable=rpm_var_1, width=12)
            rpm_entry_1.pack(side=tk.LEFT, padx=5)
            ttk.Label(rpm_frame_1, text="rpm").pack(side=tk.LEFT)
            self.entries['rpm'] = rpm_var_1
            spin_time_frame_1 = ttk.Frame(spin_frame_1)
            spin_time_frame_1.pack(fill='x', pady=3)
            ttk.Label(spin_time_frame_1, text="Time (1st):", width=15).pack(side=tk.LEFT)
            spin_time_var_1 = tk.StringVar(value=self.existing_data.get('spin_time', ''))
            spin_time_entry_1 = ttk.Entry(spin_time_frame_1, textvariable=spin_time_var_1, width=12)
            spin_time_entry_1.pack(side=tk.LEFT, padx=5)
            ttk.Label(spin_time_frame_1, text="sec").pack(side=tk.LEFT)
            self.entries['spin_time'] = spin_time_var_1
            
            # (스핀 코팅 2차)
            spin_frame_2 = ttk.LabelFrame(right_column, text="Spin Coating - 2nd Step (Optional)", padding=10)
            spin_frame_2.pack(fill='x', pady=5)
            # ... (RPM 2, Time 2 코드) ...
            rpm_frame_2 = ttk.Frame(spin_frame_2)
            rpm_frame_2.pack(fill='x', pady=3)
            ttk.Label(rpm_frame_2, text="Speed (2nd):", width=15).pack(side=tk.LEFT)
            rpm_var_2 = tk.StringVar(value=self.existing_data.get('rpm_2', ''))
            rpm_entry_2 = ttk.Entry(rpm_frame_2, textvariable=rpm_var_2, width=12)
            rpm_entry_2.pack(side=tk.LEFT, padx=5)
            ttk.Label(rpm_frame_2, text="rpm").pack(side=tk.LEFT)
            self.entries['rpm_2'] = rpm_var_2
            spin_time_frame_2 = ttk.Frame(spin_frame_2)
            spin_time_frame_2.pack(fill='x', pady=3)
            ttk.Label(spin_time_frame_2, text="Time (2nd):", width=15).pack(side=tk.LEFT)
            spin_time_var_2 = tk.StringVar(value=self.existing_data.get('spin_time_2', ''))
            spin_time_entry_2 = ttk.Entry(spin_time_frame_2, textvariable=spin_time_var_2, width=12)
            spin_time_entry_2.pack(side=tk.LEFT, padx=5)
            ttk.Label(spin_time_frame_2, text="sec").pack(side=tk.LEFT)
            self.entries['spin_time_2'] = spin_time_var_2

            # (퀜칭 방법)
            self.quenching_type_var = tk.StringVar(value=self.existing_data.get('quenching_type', 'None'))
            self.as_solvent_var = tk.StringVar(value=self.existing_data.get('as_solvent', 'Ethylacetate'))
            self.as_time_var = tk.StringVar(value=self.existing_data.get('as_time_before', ''))
            self.as_volume_var = tk.StringVar(value=self.existing_data.get('as_volume', ''))
            self.vq_duration_var = tk.StringVar(value=self.existing_data.get('vq_duration', ''))
            self.gq_start_var = tk.StringVar(value=self.existing_data.get('gq_start_before', ''))
            self.gq_duration_var = tk.StringVar(value=self.existing_data.get('gq_duration', ''))
            self.entries['quenching_type'] = self.quenching_type_var
            self.entries['as_solvent'] = self.as_solvent_var
            self.entries['as_time_before'] = self.as_time_var
            self.entries['as_volume'] = self.as_volume_var
            self.entries['vq_duration'] = self.vq_duration_var
            self.entries['gq_start_before'] = self.gq_start_var
            self.entries['gq_duration'] = self.gq_duration_var
            quenching_frame = ttk.LabelFrame(right_column, text="Quenching Method (during 1st or 2nd spin)", padding=10)
            quenching_frame.pack(fill='x', pady=5)
            type_frame = ttk.Frame(quenching_frame)
            type_frame.pack(fill='x', pady=3)
            ttk.Label(type_frame, text="Type:", width=15).pack(side=tk.LEFT)
            q_options = ["None", "Anti-Solvent", "Vacuum-Quenching", "Gas-Quenching"]
            type_combo = ttk.Combobox(
                type_frame, textvariable=self.quenching_type_var, 
                values=q_options, state='readonly', width=18
            )
            type_combo.pack(side=tk.LEFT, padx=5)
            type_combo.bind('<<ComboboxSelected>>', self._on_quenching_type_selected)
            self.quenching_detail_frame = ttk.Frame(quenching_frame, padding=(10, 5))
            self.quenching_detail_frame.pack(fill='x', expand=True)
            
            # (열처리 1차)
            anneal_frame_1 = ttk.LabelFrame(right_column, text="Annealing - 1st Step", padding=10)
            anneal_frame_1.pack(fill='x', pady=5)
            # ... (Temp 1, Time 1 코드) ...
            temp_frame_1 = ttk.Frame(anneal_frame_1)
            temp_frame_1.pack(fill='x', pady=3)
            ttk.Label(temp_frame_1, text="Temperature (1st):", width=15).pack(side=tk.LEFT)
            temp_var_1 = tk.StringVar(value=self.existing_data.get('temperature', ''))
            temp_entry_1 = ttk.Entry(temp_frame_1, textvariable=temp_var_1, width=12)
            temp_entry_1.pack(side=tk.LEFT, padx=5)
            ttk.Label(temp_frame_1, text="°C").pack(side=tk.LEFT)
            self.entries['temperature'] = temp_var_1
            anneal_time_frame_1 = ttk.Frame(anneal_frame_1)
            anneal_time_frame_1.pack(fill='x', pady=3)
            ttk.Label(anneal_time_frame_1, text="Time (1st):", width=15).pack(side=tk.LEFT)
            anneal_time_var_1 = tk.StringVar(value=self.existing_data.get('anneal_time', ''))
            anneal_time_entry_1 = ttk.Entry(anneal_time_frame_1, textvariable=anneal_time_var_1, width=12)
            anneal_time_entry_1.pack(side=tk.LEFT, padx=5)
            ttk.Label(anneal_time_frame_1, text="min").pack(side=tk.LEFT)
            self.entries['anneal_time'] = anneal_time_var_1

            # (열처리 2차)
            anneal_frame_2 = ttk.LabelFrame(right_column, text="Annealing - 2nd Step (Optional)", padding=10)
            anneal_frame_2.pack(fill='x', pady=5)
            # ... (Temp 2, Time 2 코드) ...
            temp_frame_2 = ttk.Frame(anneal_frame_2)
            temp_frame_2.pack(fill='x', pady=3)
            ttk.Label(temp_frame_2, text="Temperature (2nd):", width=15).pack(side=tk.LEFT)
            temp_var_2 = tk.StringVar(value=self.existing_data.get('temperature_2', '')) 
            temp_entry_2 = ttk.Entry(temp_frame_2, textvariable=temp_var_2, width=12)
            temp_entry_2.pack(side=tk.LEFT, padx=5)
            ttk.Label(temp_frame_2, text="°C").pack(side=tk.LEFT)
            self.entries['temperature_2'] = temp_var_2
            anneal_time_frame_2 = ttk.Frame(anneal_frame_2)
            anneal_time_frame_2.pack(fill='x', pady=3)
            ttk.Label(anneal_time_frame_2, text="Time (2nd):", width=15).pack(side=tk.LEFT)
            anneal_time_var_2 = tk.StringVar(value=self.existing_data.get('anneal_time_2', ''))
            anneal_time_entry_2 = ttk.Entry(anneal_time_frame_2, textvariable=anneal_time_var_2, width=12)
            anneal_time_entry_2.pack(side=tk.LEFT, padx=5)
            ttk.Label(anneal_time_frame_2, text="min").pack(side=tk.LEFT)
            self.entries['anneal_time_2'] = anneal_time_var_2

        # --- [공통] 추가 메모 ---
        note_frame = ttk.LabelFrame(note_frame_parent, text="Additional Notes", padding=10)
        
        if note_frame_parent == right_column:
            # 2단 레이아웃일 때: 오른쪽 하단 채우기
            note_frame.pack(fill='both', expand=True, pady=5)
        else:
            # 1단 (TCO/Contact) 레이아웃일 때: 그냥 아래에 붙이기
            note_frame.pack(fill='x', pady=5)
        
        # [신규] Placeholder 텍스트 정의
        self.placeholder_text = "#Solution\nWhen mixing DMF and DMSO at a volume ratio of 8:2,\n enter 80 and 20 respectively. (Total sum is 100) \n When using only one, enter 100.\n#Spin Coating\nIf used in 2 steps, enter twice for each step. Enter low rpm first, then high rpm.\n#Annealing\nIf the temperature changed, enter twice for each temperature.\n#Special Notes\nIf there are special notes, please write them in the Notes section."
        self.note_text = tk.Text(note_frame, height=4, width=50)
        self.note_text.pack(fill='both', expand=True)
        
        # [수정] 기존 데이터 로드 또는 Placeholder 설정
        existing_notes = self.existing_data.get('notes', '')
        if existing_notes:
            self.note_text.insert('1.0', existing_notes)
            # (기본 폰트 색상은 'black'이므로 따로 설정 안 함)
        else:
            self.note_text.insert('1.0', self.placeholder_text)
            self.note_text.config(fg='grey70') # Placeholder 색상 설정
        
        # [신규] 이벤트 바인딩
        self.note_text.bind("<FocusIn>", self._on_note_focus_in)
        self.note_text.bind("<FocusOut>", self._on_note_focus_out)
        
        # --- [공통] 하단 버튼 ---
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(side=tk.BOTTOM, pady=10) # [수정] main_frame의 맨 아래에 고정
        
        ttk.Button(btn_frame, text="Save", command=self.save_details).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_frame, text="Cancel", command=self.destroy).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_frame, text="Clear All", command=self.clear_all).pack(side=tk.LEFT, padx=5)
    
        # --- [공통] UI 초기화 ---
        if hasattr(self, 'quenching_type_var'): # 표준 UI일 때만
            self._on_quenching_type_selected() 
            self._on_deposition_method_selected()
            self._on_environment_selected()
        
    
    # --- [신규] 환경 선택 시 호출되는 함수 ---
    def _on_environment_selected(self, event=None):
        # ... (이 함수는 이전과 동일) ...
        for widget in self.env_detail_frame.winfo_children():
            widget.destroy()
        env_type = self.env_var.get()
        temp_frame = ttk.Frame(self.env_detail_frame)
        temp_frame.pack(fill='x', pady=2)
        ttk.Label(temp_frame, text="Env. Temp:", width=15).pack(side=tk.LEFT)
        temp_entry = ttk.Entry(temp_frame, textvariable=self.env_temp_var, width=12)
        temp_entry.pack(side=tk.LEFT, padx=5)
        ttk.Label(temp_frame, text="°C").pack(side=tk.LEFT)
        if env_type == "Glovebox":
            o2_frame = ttk.Frame(self.env_detail_frame)
            o2_frame.pack(fill='x', pady=2)
            ttk.Label(o2_frame, text="O₂ Level:", width=15).pack(side=tk.LEFT)
            o2_entry = ttk.Entry(o2_frame, textvariable=self.env_o2_var, width=12)
            o2_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(o2_frame, text="ppm").pack(side=tk.LEFT)
            h2o_frame = ttk.Frame(self.env_detail_frame)
            h2o_frame.pack(fill='x', pady=2)
            ttk.Label(h2o_frame, text="H₂O Level:", width=15).pack(side=tk.LEFT)
            h2o_entry = ttk.Entry(h2o_frame, textvariable=self.env_h2o_var, width=12)
            h2o_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(h2o_frame, text="ppm").pack(side=tk.LEFT)
        elif env_type == "Air":
            hum_frame = ttk.Frame(self.env_detail_frame)
            hum_frame.pack(fill='x', pady=2)
            ttk.Label(hum_frame, text="Env. Humidity:", width=15).pack(side=tk.LEFT)
            hum_entry = ttk.Entry(hum_frame, textvariable=self.env_humidity_var, width=12)
            hum_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(hum_frame, text="% RH").pack(side=tk.LEFT)

    # --- 증착 방법 선택 시 호출되는 함수 ---
    def _on_deposition_method_selected(self, event=None):
        # ... (이 함수는 이전과 동일) ...
        for widget in self.dep_method_detail_frame.winfo_children():
            widget.destroy()
        if self.method_var.get() == "Dynamic":
            dyn_time_frame = ttk.Frame(self.dep_method_detail_frame)
            dyn_time_frame.pack(fill='x', pady=2)
            ttk.Label(dyn_time_frame, text="Drop Time:", width=15).pack(side=tk.LEFT)
            dyn_time_entry = ttk.Entry(dyn_time_frame, textvariable=self.dynamic_drop_time_var, width=12)
            dyn_time_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(dyn_time_frame, text="sec (before finished)").pack(side=tk.LEFT)

    # --- [신규] Note Text Placeholder 헬퍼 함수 ---
    def _on_note_focus_in(self, event=None):
        """노트 텍스트 박스를 클릭했을 때"""
        current_text = self.note_text.get('1.0', 'end-1c').strip()
        if current_text == self.placeholder_text:
            self.note_text.delete('1.0', 'end')
            self.note_text.config(fg='black') # (시스템 기본값으로 변경)

    def _on_note_focus_out(self, event=None):
        """노트 텍스트 박스에서 포커스가 나갔을 때"""
        current_text = self.note_text.get('1.0', 'end-1c').strip()
        if not current_text:
            self.note_text.insert('1.0', self.placeholder_text)
            self.note_text.config(fg='grey70')
    # --- [신규 함수 끝] ---

    # --- 퀜칭 타입 선택 시 호출되는 함수 ---
    def _on_quenching_type_selected(self, event=None):
        # ... (이 함수는 이전과 동일, 오타 수정: width=1V -> width=12) ...
        for widget in self.quenching_detail_frame.winfo_children():
            widget.destroy()
        q_type = self.quenching_type_var.get()
        if q_type == "Anti-Solvent":
            as_solvent_frame = ttk.Frame(self.quenching_detail_frame)
            as_solvent_frame.pack(fill='x', pady=2)
            ttk.Label(as_solvent_frame, text="Solvent:", width=15).pack(side=tk.LEFT)
            as_options = ["Ethylacetate", "Diethyl ether", "Chlorobenzene", "Toluene", "Anisole"]
            as_combo = ttk.Combobox(
                as_solvent_frame, textvariable=self.as_solvent_var,
                values=as_options, width=18
            )
            as_combo.pack(side=tk.LEFT, padx=5)
            as_time_frame = ttk.Frame(self.quenching_detail_frame)
            as_time_frame.pack(fill='x', pady=2)
            ttk.Label(as_time_frame, text="Drop Time:", width=15).pack(side=tk.LEFT)
            as_time_entry = ttk.Entry(as_time_frame, textvariable=self.as_time_var, width=12)
            as_time_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(as_time_frame, text="sec (before finished)").pack(side=tk.LEFT)
            as_volume_frame = ttk.Frame(self.quenching_detail_frame)
            as_volume_frame.pack(fill='x', pady=2)
            ttk.Label(as_volume_frame, text="Drop Volume:", width=15).pack(side=tk.LEFT)
            as_volume_entry = ttk.Entry(as_volume_frame, textvariable=self.as_volume_var, width=12) # [오타 수정] 1V -> 12
            as_volume_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(as_volume_frame, text="μL").pack(side=tk.LEFT)
        elif q_type == "Vacuum-Quenching":
            vq_frame = ttk.Frame(self.quenching_detail_frame)
            vq_frame.pack(fill='x', pady=2)
            ttk.Label(vq_frame, text="Duration:", width=15).pack(side=tk.LEFT)
            vq_entry = ttk.Entry(vq_frame, textvariable=self.vq_duration_var, width=12)
            vq_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(vq_frame, text="sec").pack(side=tk.LEFT)
        elif q_type == "Gas-Quenching":
            gq_start_frame = ttk.Frame(self.quenching_detail_frame)
            gq_start_frame.pack(fill='x', pady=2)
            ttk.Label(gq_start_frame, text="Start Time:", width=15).pack(side=tk.LEFT)
            gq_start_entry = ttk.Entry(gq_start_frame, textvariable=self.gq_start_var, width=12)
            gq_start_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(gq_start_frame, text="sec (before finished)").pack(side=tk.LEFT)
            gq_duration_frame = ttk.Frame(self.quenching_detail_frame)
            gq_duration_frame.pack(fill='x', pady=2)
            ttk.Label(gq_duration_frame, text="Duration:", width=15).pack(side=tk.LEFT)
            gq_duration_entry = ttk.Entry(gq_duration_frame, textvariable=self.gq_duration_var, width=12)
            gq_duration_entry.pack(side=tk.LEFT, padx=5)
            ttk.Label(gq_duration_frame, text="sec").pack(side=tk.LEFT)

    
    def save_details(self):
        """[수정됨] 입력된 모든 상세정보(환경, 용매 비율 포함)를 저장합니다."""
        details = {}
        for key, var in self.entries.items():
            value = var.get().strip()
            if value:
                details[key] = value
        
            # [수정] Placeholder 텍스트는 저장하지 않도록 처리
        notes = self.note_text.get('1.0', 'end-1c').strip()
        if notes and notes != self.placeholder_text:
            details['notes'] = notes
        
        if hasattr(self, 'quenching_type_var'): # 표준 UI일 때만
            # 'Static'일 때 불필요한 dynamic_drop_time 키 삭제
            if details.get('deposition_method') == 'Static':
                details.pop('dynamic_drop_time', None)
            
            # 'None' 퀜칭일 때 불필요한 퀜칭 키 삭제
            if details.get('quenching_type') == 'None':
                keys_to_pop = [k for k in details if k.startswith('as_') or k.startswith('vq_') or k.startswith('gq_')]
                for k in keys_to_pop:
                    details.pop(k, None)
            
            # [신규] 환경에 따라 불필요한 키 삭제
            env_type = details.get('spin_environment')
            if env_type == 'Glovebox':
                details.pop('spin_env_humidity', None) # 습도 키 삭제
            elif env_type == 'Air':
                details.pop('spin_env_o2', None) # O2 키 삭제
                details.pop('spin_env_h2o', None) # H2O 키 삭제
        
        process_details[self.sample_name][self.variable_name][self.material] = details
        
        messagebox.showinfo("Saved", f"Process details for {self.material} have been saved.")
        self.destroy()
    
    def clear_all(self):
        """[수정됨] 모든 입력 필드를 초기화합니다."""
        if messagebox.askyesno("Clear All", "Are you sure you want to clear all fields?"):
            for var in self.entries.values():
                var.set('')
                self.note_text.delete('1.0', 'end')
                self.note_text.insert('1.0', self.placeholder_text)
                self.note_text.config(fg='grey70')
            
            if hasattr(self, 'quenching_type_var'): # 표준 UI일 때만
                self.quenching_type_var.set('None')
                self._on_quenching_type_selected()
                
                self.method_var.set('Static')
                self._on_deposition_method_selected()
                
                # [신규] 환경 UI 'Glovebox'로 리셋
                self.env_var.set('Glovebox')
                self._on_environment_selected()

# --- 다중 선택 팝업창을 여는 헬퍼 함수 ---
def open_variable_selector(col_name, result_var):
    """다중 선택 팝업창을 여는 헬퍼 함수 (공정 조건 입력 기능 추가)"""
    options = VARIABLE_PRESETS.get(col_name, [])
    # [수정] 옵션이 없어도 커스텀 입력을 위해 창을 열도록 수정
    # if not options:
    #     messagebox.showinfo("Info", f"No presets defined for '{col_name}'.")
    #     return
    EnhancedVariableSelectorWindow(root, f"Select {col_name}", col_name, options, result_var)


# --- Process Details Viewer Window ---
class ProcessDetailsViewerWindow(tk.Toplevel):
    def __init__(self, parent, sample_name):
        super().__init__(parent)
        self.title(f"Process Details for {sample_name}")
        self.geometry("600x900")
        self.transient(parent)
        self.grab_set()

        self.sample_name = sample_name

        main_frame = ttk.Frame(self, padding=10)
        main_frame.pack(fill='both', expand=True)

        # Display details here
        details_text = tk.Text(main_frame, wrap='word', state='disabled')
        details_text.pack(fill='both', expand=True)

        # Populate text widget with process_details
        if sample_name in process_details:
            all_details = process_details[sample_name]
            formatted_details = []
            for var_name, materials in all_details.items():
                formatted_details.append(f"--- {var_name} ---")
                for material, data in materials.items():
                    formatted_details.append(f"  Material: {material}")
                    for key, value in data.items():
                        formatted_details.append(f"    {key.replace('_', ' ').capitalize()}: {value}")
                formatted_details.append("\n")

            details_text.config(state='normal')
            details_text.insert('1.0', "\n".join(formatted_details))
            details_text.config(state='disabled')
        else:
            details_text.config(state='normal')
            details_text.insert('1.0', "No process details available for this sample.")
            details_text.config(state='disabled')

        ttk.Button(main_frame, text="Close", command=self.destroy).pack(pady=10)

# --- Function to open the Process Details Viewer ---
def view_process_details():
    sample_name = current_sample_label.get()
    if sample_name == "[No sample selected]":
        messagebox.showwarning("No Sample", "Please select a sample folder from the list first to view process details.")
        return
    ProcessDetailsViewerWindow(root, sample_name)

# --- Variable Manager Window ---
class VariableManagerWindow(tk.Toplevel):
    def __init__(self, parent):
        super().__init__(parent)
        self.title("Experimental Variable Manager")
        self.geometry("900x400")
        
        self.samples = sorted(current_display_df['Sample'].unique())
        self.var_data = experimental_variables
        
        cols = ["Sample"] + variable_columns
        
        frame = ttk.Frame(self)
        frame.pack(fill='both', expand=True, padx=10, pady=10)
        
        self.var_tree = ttk.Treeview(frame, columns=cols, show='headings')
        
        vsb = ttk.Scrollbar(frame, orient="vertical", command=self.var_tree.yview)
        vsb.pack(side='right', fill='y')
        hsb = ttk.Scrollbar(frame, orient="horizontal", command=self.var_tree.xview)
        hsb.pack(side='bottom', fill='x')
        self.var_tree.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set)
        
        for col in cols:
            self.var_tree.heading(col, text=col)
            self.var_tree.column(col, width=100, anchor='center')
        self.var_tree.column("Sample", anchor='w', width=150, minwidth=100)
        
        for sample in self.samples:
            values = [sample] + [self.var_data.get(sample, {}).get(col, "") for col in variable_columns]
            self.var_tree.insert('', 'end', values=values)
            
        self.var_tree.pack(fill='both', expand=True)
        self.var_tree.bind("<Double-1>", self.on_double_click)
        
        ttk.Button(self, text="Save and Close", command=self.save_and_close).pack(pady=10)

    def on_double_click(self, event):
        region = self.var_tree.identify("region", event.x, event.y)
        if region != "cell": 
            return

        col_id = self.var_tree.identify_column(event.x)
        col_index = int(col_id.replace("#", "")) - 1

        if col_index == 0: 
            return

        item_id = self.var_tree.identify_row(event.y)
        x, y, width, height = self.var_tree.bbox(item_id, col_id)

        entry_var = tk.StringVar()
        entry = ttk.Entry(self.var_tree, textvariable=entry_var)
        entry.place(x=x, y=y, width=width, height=height)

        current_val = self.var_tree.item(item_id, 'values')[col_index]
        entry_var.set(current_val)
        entry.focus_force()

        def save_edit(event):
            new_val = entry_var.get()
            current_values = list(self.var_tree.item(item_id, 'values'))
            current_values[col_index] = new_val
            self.var_tree.item(item_id, values=current_values)
            entry.destroy()

        entry.bind("<Return>", save_edit)
        entry.bind("<FocusOut>", save_edit)

    def save_and_close(self):
        global experimental_variables, current_display_df
        temp_experimental_variables = {}

        for item_id in self.var_tree.get_children():
            values = self.var_tree.item(item_id, 'values')
            sample_name = values[0]
            temp_experimental_variables[sample_name] = {}
            for i, col in enumerate(variable_columns):
                temp_experimental_variables[sample_name][col] = values[i+1]
                is_var_key = f"{col}_is_var"
                if sample_name in experimental_variables and is_var_key in experimental_variables[sample_name]:
                    temp_experimental_variables[sample_name][is_var_key] = experimental_variables[sample_name][is_var_key]
                elif sample_name not in experimental_variables:
                    temp_experimental_variables[sample_name][is_var_key] = False

        experimental_variables = temp_experimental_variables

        merged_vars_df = pd.DataFrame.from_dict(experimental_variables, orient='index').reset_index().rename(columns={'index': 'Sample'})

        all_var_cols_to_consider = variable_columns + [f"{col}_is_var" for col in variable_columns]
        cols_to_drop = [
            col for col in all_var_cols_to_consider 
            if col in current_display_df.columns and col != 'Sample'
        ]

        df_without_vars = current_display_df.drop(columns=cols_to_drop, errors='ignore')
        current_display_df = pd.merge(df_without_vars, merged_vars_df, on="Sample", how="left")

        refresh_all_views(current_display_df)
        self.destroy()


# --- [NEW] Unified Variable Manager Window ---
class UnifiedVariableManagerWindow(tk.Toplevel):
    """통합 변수 관리 창 - 모든 변수를 한 화면에서 입력"""
    
    def __init__(self, parent, sample_name):
        super().__init__(parent)
        self.title(f"Unified Variable Manager - {sample_name}")
        self.geometry("1000x700")
        self.transient(parent)
        self.grab_set()
        
        self.sample_name = sample_name
        self.parent = parent
        
        # 각 변수별 입력 위젯을 저장할 딕셔너리
        self.variable_widgets = {}
        self.detail_frames = {}  # 접기/펼치기를 위한 프레임 저장
        
        # 메인 컨테이너
        main_container = ttk.Frame(self)
        main_container.pack(fill='both', expand=True)
        
        # 헤더
        header_frame = ttk.Frame(main_container)
        header_frame.pack(fill='x', padx=10, pady=10)
        ttk.Label(header_frame, text=f"Sample: {sample_name}", 
                 font=('Helvetica', 12, 'bold')).pack(side=tk.LEFT)
        
        # 스크롤 가능한 캔버스 생성
        canvas = tk.Canvas(main_container)
        scrollbar = ttk.Scrollbar(main_container, orient="vertical", command=canvas.yview)
        self.scrollable_frame = ttk.Frame(canvas)
        
        self.scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=self.scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        canvas.pack(side="left", fill="both", expand=True, padx=(10, 0), pady=5)
        scrollbar.pack(side="right", fill="y", pady=5, padx=(0, 10))
        
        # 마우스 휠 스크롤 바인딩
        def _on_mousewheel(event):
            canvas.yview_scroll(int(-1*(event.delta/120)), "units")
        canvas.bind_all("<MouseWheel>", _on_mousewheel)
        
        # 각 변수별 섹션 생성
        self._create_variable_sections()
        
        # 하단 버튼
        button_frame = ttk.Frame(main_container)
        button_frame.pack(fill='x', padx=10, pady=10)
        
        ttk.Button(button_frame, text="Save All Variables", 
                  command=self.save_all_variables).pack(side=tk.LEFT, padx=5)
        ttk.Button(button_frame, text="Load Recipe...", 
                  command=self.load_recipe_into_window).pack(side=tk.LEFT, padx=5)
        ttk.Button(button_frame, text="Cancel", 
                  command=self.destroy).pack(side=tk.RIGHT, padx=5)
        
        # 창이 닫힐 때 마우스 휠 바인딩 해제
        self.protocol("WM_DELETE_WINDOW", lambda: (canvas.unbind_all("<MouseWheel>"), self.destroy()))
    
    def _create_variable_sections(self):
        """각 변수별 섹션을 생성합니다."""
        
        # 기존 데이터 로드
        sample_vars = experimental_variables.get(self.sample_name, {})
        sample_details = process_details.get(self.sample_name, {})
        
        for col in variable_columns:
            # 변수별 LabelFrame 생성
            var_frame = ttk.LabelFrame(self.scrollable_frame, text=col, padding=10)
            var_frame.pack(fill='x', padx=5, pady=5)
            
            # 재료 선택 프레임
            material_frame = ttk.Frame(var_frame)
            material_frame.pack(fill='x', pady=(0, 5))
            
            ttk.Label(material_frame, text="Material:", width=10).pack(side=tk.LEFT)
            
            # 재료 입력 변수
            material_var = tk.StringVar(value=sample_vars.get(col, ""))
            
            if col == "Perovskite":
                # Perovskite는 readonly + Build 버튼
                entry = ttk.Entry(material_frame, textvariable=material_var, 
                                 state='readonly', width=40)
                entry.pack(side=tk.LEFT, padx=5)
                
                ttk.Button(material_frame, text="Build...", width=8,
                          command=lambda v=material_var: self._open_perovskite_builder(v)).pack(side=tk.LEFT, padx=2)
            elif col in VARIABLE_PRESETS:
                # 프리셋이 있는 변수는 readonly + ... 버튼
                entry = ttk.Entry(material_frame, textvariable=material_var, 
                                 state='readonly', width=40)
                entry.pack(side=tk.LEFT, padx=5)
                
                ttk.Button(material_frame, text="...", width=5,
                          command=lambda c=col, v=material_var: self._open_material_selector(c, v)).pack(side=tk.LEFT, padx=2)
            else:
                # 일반 입력
                entry = ttk.Entry(material_frame, textvariable=material_var, width=40)
                entry.pack(side=tk.LEFT, padx=5)
            
            self.variable_widgets[col] = material_var
            
            # 공정 조건 프레임 (접기/펼치기 가능)
            if col not in ["Perovskite"]:  # Perovskite는 공정 조건 없음
                detail_container = ttk.Frame(var_frame)
                detail_container.pack(fill='x', pady=(5, 0))
                
                # 접기/펼치기 버튼
                toggle_var = tk.BooleanVar(value=False)
                toggle_btn = ttk.Checkbutton(
                    detail_container, 
                    text="▶ Process Details",
                    variable=toggle_var,
                    command=lambda c=col, tv=toggle_var, dc=detail_container: self._toggle_details(c, tv, dc)
                )
                toggle_btn.pack(anchor='w')
                
                # 상세 정보 프레임 (초기에는 숨김)
                detail_frame = ttk.Frame(detail_container)
                self.detail_frames[col] = {
                    'frame': detail_frame,
                    'toggle_var': toggle_var,
                    'toggle_btn': toggle_btn,
                    'container': detail_container
                }
                
                # 기존 공정 조건이 있으면 표시
                if col in sample_details and sample_details[col]:
                    self._create_detail_summary(detail_container, col, sample_details[col])
    
    def _toggle_details(self, col, toggle_var, container):
        """공정 조건 섹션을 접거나 펼칩니다."""
        detail_info = self.detail_frames[col]
        detail_frame = detail_info['frame']
        toggle_btn = detail_info['toggle_btn']
        
        if toggle_var.get():
            # 펼치기
            toggle_btn.config(text="▼ Process Details")
            detail_frame.pack(fill='x', padx=20, pady=5)
            
            # 공정 조건 입력 필드 생성 (간소화 버전)
            self._create_simplified_detail_inputs(detail_frame, col)
        else:
            # 접기
            toggle_btn.config(text="▶ Process Details")
            detail_frame.pack_forget()
    
    def _create_detail_summary(self, container, col, materials_dict):
        """기존 공정 조건 요약 표시"""
        summary_frame = ttk.Frame(container)
        summary_frame.pack(fill='x', padx=20, pady=2)
        
        summary_text = []
        for material, details in materials_dict.items():
            if details:
                summary_text.append(f"  • {material}: {len(details)} parameters")
        
        if summary_text:
            ttk.Label(summary_frame, text="\n".join(summary_text), 
                     foreground='gray').pack(anchor='w')
    
    def _create_simplified_detail_inputs(self, frame, col):
        """간소화된 공정 조건 입력 필드 생성"""
        # 기존 위젯 제거
        for widget in frame.winfo_children():
            widget.destroy()
        
        info_label = ttk.Label(frame, 
                              text="Click '...' button to edit detailed process parameters",
                              foreground='blue', cursor='hand2')
        info_label.pack(anchor='w', pady=5)
        
        # 클릭 시 상세 입력 창 열기
        info_label.bind("<Button-1>", 
                       lambda e, c=col: self._open_detail_editor(c))
    
    def _open_detail_editor(self, col):
        """상세 공정 조건 편집 창 열기"""
        material_str = self.variable_widgets[col].get()
        if not material_str:
            messagebox.showinfo("No Material", 
                              f"Please select a material for {col} first.",
                              parent=self)
            return
        
        # 첫 번째 재료에 대한 상세 정보 창 열기
        materials = [m.strip() for m in material_str.split(' + ') if m.strip()]
        if materials:
            ProcessDetailWindow(self, self.sample_name, col, materials[0])
    
    def _open_perovskite_builder(self, result_var):
        """Perovskite Builder 창 열기"""
        PerovskiteBuilderWindow(self, self.sample_name)
        # 빌더 창이 닫힌 후 값 업데이트는 자동으로 처리됨
    
    def _open_material_selector(self, col, result_var):
        """재료 선택 창 열기"""
        options = VARIABLE_PRESETS.get(col, [])
        EnhancedVariableSelectorWindow(self, f"Select {col}", col, options, result_var)
    
    def load_recipe_into_window(self):
        """레시피를 불러와서 현재 창에 적용"""
        filepath = filedialog.askopenfilename(
            title="Load Variable Recipe",
            filetypes=[("JV Recipe Files", "*.jvr"), ("All files", "*.*")],
            parent=self
        )
        if not filepath:
            return
        
        try:
            with open(filepath, 'rb') as f:
                recipe_data = pickle.load(f)
            
            if not isinstance(recipe_data, dict) or 'variables' not in recipe_data:
                messagebox.showerror("Load Error", "Invalid recipe file.", parent=self)
                return
            
            loaded_vars = recipe_data.get('variables', {})
            loaded_details = recipe_data.get('details', {})
            
            # GUI에 값 적용
            for col in variable_columns:
                if col in loaded_vars and col in self.variable_widgets:
                    self.variable_widgets[col].set(loaded_vars[col])
            
            # 공정 상세정보도 메모리에 임시 적용
            if self.sample_name not in process_details:
                process_details[self.sample_name] = {}
            
            for col, materials_dict in loaded_details.items():
                process_details[self.sample_name][col] = materials_dict.copy()
            
            messagebox.showinfo("Recipe Loaded", 
                              "Recipe has been loaded into the form.\n"
                              "Click 'Save All Variables' to apply.",
                              parent=self)
        
        except Exception as e:
            messagebox.showerror("Load Error", f"Failed to load recipe:\n{e}", parent=self)
    
    def save_all_variables(self):
        """모든 변수를 저장"""
        global experimental_variables, current_display_df
        
        # experimental_variables 업데이트
        if self.sample_name not in experimental_variables:
            experimental_variables[self.sample_name] = {}
        
        control_vars_to_fill = {}
        
        for col in variable_columns:
            value = self.variable_widgets[col].get()
            experimental_variables[self.sample_name][col] = value
            
            # is_var 플래그는 기존 값 유지
            is_var_key = f"{col}_is_var"
            if is_var_key not in experimental_variables[self.sample_name]:
                experimental_variables[self.sample_name][is_var_key] = False
            
            # Control 변수 수집 (체크되지 않은 변수)
            is_variable = experimental_variables[self.sample_name].get(is_var_key, False)
            if not is_variable and value:
                control_vars_to_fill[col] = value
        
        # Control 변수를 다른 샘플에도 적용할지 물어보기
        if control_vars_to_fill:
            msg = "The following control variables will be applied to ALL other samples:\n\n"
            msg += "\n".join([f"{col}: {val}" for col, val in control_vars_to_fill.items()])
            msg += "\n\nProceed?"
            
            if messagebox.askyesno("Apply Control Variables?", msg, parent=self):
                all_samples = current_display_df['Sample'].unique()
                for s in all_samples:
                    if s == self.sample_name:
                        continue
                    if s not in experimental_variables:
                        experimental_variables[s] = {}
                    
                    for col, val in control_vars_to_fill.items():
                        experimental_variables[s][col] = val
                        
                        # 공정 상세정보도 복사
                        materials_to_copy = [m.strip() for m in val.split(' + ') if m.strip()]
                        for material in materials_to_copy:
                            source_details = process_details.get(self.sample_name, {}).get(col, {}).get(material, {})
                            if source_details:
                                if s not in process_details:
                                    process_details[s] = {}
                                if col not in process_details[s]:
                                    process_details[s][col] = {}
                                process_details[s][col][material] = source_details.copy()
        
        # DataFrame 업데이트
        merged_vars_df = pd.DataFrame.from_dict(
            experimental_variables, 
            orient='index'
        ).reset_index().rename(columns={'index': 'Sample'})
        
        all_var_cols_to_consider = variable_columns + [f"{col}_is_var" for col in variable_columns]
        cols_to_drop = [
            col for col in all_var_cols_to_consider 
            if col in current_display_df.columns and col != 'Sample'
        ]
        
        df_without_vars = current_display_df.drop(columns=cols_to_drop, errors='ignore')
        current_display_df = pd.merge(df_without_vars, merged_vars_df, on="Sample", how="left")
        
        refresh_all_views(current_display_df)
        
        messagebox.showinfo("Saved", 
                          f"All variables for '{self.sample_name}' have been saved.",
                          parent=self)
        self.destroy()



def open_variable_manager():
    """통합 변수 관리 창을 엽니다."""
    sample_name = current_sample_label.get()
    if sample_name == "[No sample selected]":
        messagebox.showwarning("No Sample", "Please select a sample from the list first.")
        return
    
    UnifiedVariableManagerWindow(root, sample_name)

# --- NEW: PowerPoint Export Function ---
def export_to_powerpoint():
    if not plotted_jv_items and not plotted_dist_items and pce_table_best.get_children() == ():
        messagebox.showwarning("No Data", "Please plot J-V curves, distribution graphs, or have data in the 'Best Device' table to export.")
        return

    filepath = filedialog.asksaveasfilename(
        title="Save PowerPoint Report",
        defaultextension=".pptx",
        filetypes=[("PowerPoint Presentations", "*.pptx"), ("All files", "*.*")]
    )
    if not filepath:
        return

    try:
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        if plotted_jv_items:
            img_jv_stream = io.BytesIO()
            fig_jv.savefig(img_jv_stream, format='png', dpi=300, bbox_inches='tight')
            img_jv_stream.seek(0)
            slide.shapes.add_picture(img_jv_stream, Inches(0.5), Inches(0.5), width=Inches(4.5))
            img_jv_stream.close()

        if plotted_dist_items:
            img_dist_stream = io.BytesIO()
            fig_dist.savefig(img_dist_stream, format='png', dpi=300, bbox_inches='tight')
            img_dist_stream.seek(0)
            slide.shapes.add_picture(img_dist_stream, Inches(5.0), Inches(0.5), width=Inches(4.5))
            img_dist_stream.close()

        table_items = pce_table_best.get_children()
        if table_items:
            cols = pce_table_best['columns']
            rows = len(table_items) + 1 

            table_shape = slide.shapes.add_table(rows, len(cols), Inches(0.5), Inches(4.5), Inches(9.0), Inches(0.5)) 
            table = table_shape.table

            for c_idx, col_name in enumerate(cols):
                cell = table.cell(0, c_idx)
                cell.text = col_name
                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.size = Pt(9)
                para.alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220, 220, 220) 

            for r_idx, item_id in enumerate(table_items):
                values = pce_table_best.item(item_id)['values']
                for c_idx, value in enumerate(values):
                    cell = table.cell(r_idx + 1, c_idx)
                    cell.text = str(value)
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(8)
                    para.alignment = PP_ALIGN.CENTER if c_idx > 1 else PP_ALIGN.LEFT 

            try:
                table.columns[0].width = Inches(1.5) # Sample
                table.columns[1].width = Inches(2.0) # File
            except IndexError:
                 pass 

            for row in table.rows:
                row.height = Pt(15)

        prs.save(filepath)
        messagebox.showinfo("Export Complete", f"Report saved successfully to:\n{filepath}")

    except Exception as e:
        messagebox.showerror("Export Error", f"Failed to export to PowerPoint:\n{e}")

# --- [신규] Deep Learning / ML용 데이터셋 내보내기 함수 (파일명 자동 생성) ---
def export_for_ml():
    # [수정] operator_name, device_structure 전역 변수 추가
    global current_display_df, experimental_variables, process_details, variable_columns, \
           current_root_folder, var_checkbox_vars, \
           operator_name, device_structure
    
    if current_display_df.empty:
        messagebox.showwarning("No Data", "Please load data first.", parent=root)
        return

    # --- [수정] 파일명 자동 생성 로직 (사용자 이름, 소자 구조 포함) ---
    try:
        # 1. 폴더 날짜 (e.g., "20251110")
        folder_date = "UnknownDate"
        if current_root_folder:
            # os.path.basename()을 사용하여 폴더 경로의 마지막 이름만 가져옵니다.
            folder_date = os.path.basename(current_root_folder)

        # 2. 체크된 변수 (e.g., "HTL_Perovskite")
        # 'Experimental Variables' 패널에서 체크된 변수들의 리스트를 가져옵니다.
        varying_vars = [col for col in variable_columns if col in var_checkbox_vars and var_checkbox_vars[col].get()]
        
        if not varying_vars:
            vars_str = "Batch" # 아무것도 체크되지 않았을 때의 기본값
        else:
            vars_str = "_".join(varying_vars) # 예: "HTL_Perovskite"

        # 3. 밴드갭 (e.g., "1-67eV")
        bandgap_str = "NA_Eg" # 기본값
        target_samples = current_display_df['Sample'].unique()
        
        # 현재 표시된 샘플들을 순회하며 'Perovskite' 변수에서 밴드갭(Eg)을 찾습니다.
        for sample in target_samples:
            if sample in experimental_variables and "Perovskite" in experimental_variables[sample]:
                p_string = experimental_variables[sample]["Perovskite"]
                # 정규표현식(re)을 사용해 "(Eg=1.67eV)" 같은 패턴을 찾습니다.
                match = re.search(r"\(Eg=([\d\.]+)\s*eV\)", p_string)
                if match:
                    # 찾은 밴드갭(예: "1.67")의 '.'을 '-'로 바꿉니다. (예: "1-67eV")
                    bandgap_str = match.group(1).replace(".", "-") + "eV"
                    break # 첫 번째 밴드갭을 찾으면 중단

        # 4. [수정] 사용자가 직접 입력할 필드 (구조, 사용자명)
        # 전역 변수에서 값을 가져오고, 없으면 기본 문자열 사용
        structure_str = device_structure if device_structure else "[Structure]"
        operator_str = operator_name if operator_name else "[UserName]"
        user_fields = f"{structure_str}_{operator_str}"
        
        # 5. 모든 조각을 조합하여 최종 파일명 제안
        initial_filename = f"{folder_date}_{vars_str}_{bandgap_str}_{user_fields}.csv"

    except Exception as e:
        print(f"Error generating filename: {e}")
        initial_filename = "ML_Export.csv"
    # --- [수정된 로직 끝] ---


    filepath = filedialog.asksaveasfilename(
        title="Export for Deep Learning",
        initialfile=initial_filename, # [수정] 자동 생성된 파일명을 기본값으로 제안
        defaultextension=".csv",
        filetypes=[("CSV (Comma-separated values)", "*.csv"), ("Excel Files", "*.xlsx"), ("All files", "*.*")],
        parent=root
    )
    if not filepath:
        return

    #
    # (이하 데이터 수집 및 저장 로직은 이전과 동일합니다)
    #
    all_devices_data_for_export = []
    
    for i in range(len(current_display_df)):
        device_row = current_display_df.iloc[i]
        sample_name = device_row.get('Sample')
        if not sample_name:
            continue

        device_data_dict = {}

        # --- [신규] 1. 배치 정보 (사용자, 구조) 추가 ---
        device_data_dict['Operator'] = operator_name
        device_data_dict['Structure'] = device_structure
        
        # --- 2. 장치 기본 정보 및 타겟 변수(PCE 등) 추가 ---
        device_data_dict['Sample'] = sample_name
        device_data_dict['File'] = device_row.get('File')
        device_data_dict['Scan Direction'] = device_row.get('Scan Direction')
        device_data_dict['PCE (%)'] = device_row.get('PCE (%)')
        device_data_dict['Voc (V)'] = device_row.get('Voc (V)')
        device_data_dict['Jsc (mA/cm2)'] = device_row.get('Jsc (mA/cm2)')
        device_data_dict['FF (%)'] = device_row.get('FF (%)')
        device_data_dict['Rs (Ω·cm²)'] = device_row.get('Rs (Ω·cm²)')
        device_data_dict['Rsh (Ω·cm²)'] = device_row.get('Rsh (Ω·cm²)')

        # --- 3. 레벨 1 변수 (7가지 주요 변수 문자열) 추가 ---
        sample_vars = experimental_variables.get(sample_name, {})
        for col_name in variable_columns: # TCO, HTL, Perovskite...
            device_data_dict[col_name] = sample_vars.get(col_name, np.nan) # 없으면 NaN

        # --- 4. 레벨 2 변수 (모든 상세 공정 파라미터) 추가 (수정됨) ---
        sample_details = process_details.get(sample_name, {})
        
        for var_name, materials_dict in sample_details.items():
            materials_string = sample_vars.get(var_name, "")
            if pd.isna(materials_string):
                materials_string = ""
            current_materials_list = [m.strip() for m in materials_string.split(' + ') if m.strip()]

            for material_name, params_dict in materials_dict.items():
                if material_name not in current_materials_list:
                    continue # "고아" 데이터 무시

                clean_material_name = re.sub(r'[^A-Za-z0-9_]', '_', material_name)
                
                for param_name, param_value in params_dict.items():
                    flat_col_name = f"{var_name}_{clean_material_name}_{param_name}"
                    
                    try:
                        numeric_value = pd.to_numeric(param_value)
                    except (ValueError, TypeError):
                        numeric_value = param_value
                    
                    device_data_dict[flat_col_name] = numeric_value
        
        all_devices_data_for_export.append(device_data_dict)

    if not all_devices_data_for_export:
        messagebox.showwarning("No Data", "분석할 데이터가 없습니다.", parent=root)
        return

    # --- 5. 딕셔너리 리스트를 Pandas DataFrame으로 변환 ---
    export_df = pd.DataFrame(all_devices_data_for_export)

    # --- 6. XGBoost 호환을 위해 컬럼명 정제 (특수문자 제거) ---
    export_df = clean_column_names_for_ml(export_df)

    # --- 7. 파일로 저장 ---
    try:
        if filepath.endswith('.csv'):
            export_df.to_csv(filepath, index=False, encoding='utf-8-sig')
        elif filepath.endswith('.xlsx'):
            export_df.to_excel(filepath, index=False)
        
        messagebox.showinfo("Export Complete", 
                            f"{len(export_df)}개 장치의 데이터가 ML용으로 저장되었습니다:\n{filepath}\n\n"
                            f"컬럼명이 XGBoost 호환 형식으로 정제되었습니다.\n"
                            f"(특수문자 제거: (, ), [, ], <, >, · 등)",
                            parent=root)
    except Exception as e:
        messagebox.showerror("Export Error", f"파일 저장에 실패했습니다:\n{e}", parent=root)

# --- GUI 창 생성 및 레이아웃 설정 ---
root = ThemedTk(theme="arc")
root.title("Solar Cell J-V & Statistics Analyzer for LKChem")
root.state('zoomed')
root.protocol("WM_DELETE_WINDOW", on_closing)

# --- Menu Bar ---
menubar = tk.Menu(root)
root.config(menu=menubar)
file_menu = tk.Menu(menubar, tearoff=0)
file_menu.add_command(label="Open Folder...", command=load_and_process_folder)
file_menu.add_separator()
file_menu.add_command(label="Save Analysis State...", command=save_state)
file_menu.add_command(label="Load Analysis State...", command=load_state)
file_menu.add_separator()
file_menu.add_command(label="Save Variable Recipe...", command=save_variable_recipe)
file_menu.add_command(label="Load Variable Recipe...", command=load_variable_recipe)
file_menu.add_separator()
file_menu.add_command(label="Export to PowerPoint...", command=export_to_powerpoint)
file_menu.add_command(label="Export for ML (.csv/xlsx)...", command=export_for_ml)
file_menu.add_separator()
file_menu.add_command(label="Exit", command=on_closing)
menubar.add_cascade(label="File", menu=file_menu)

edit_menu = tk.Menu(menubar, tearoff=0)
edit_menu.add_command(label="Auto Merge Folders by Number", command=auto_merge_folders)
edit_menu.add_command(label="Best PCE (Sep. Fwd/Rev)", command=filter_best_pce_per_sample_direction)
menubar.add_cascade(label="Edit", menu=edit_menu)



# --- Main Layout ---
main_paned_window = ttk.PanedWindow(root, orient=tk.VERTICAL)
main_paned_window.pack(fill=tk.BOTH, expand=True)
top_frame = ttk.Frame(main_paned_window)
bottom_frame = ttk.Frame(main_paned_window)
main_paned_window.add(top_frame)
main_paned_window.add(bottom_frame)

def adjust_sash(event=None): # [수정 1] (event=None) 추가
    main_paned_window.update_idletasks()
    try:
        # 1. 창의 전체 세로 높이를 가져옵니다.
        total_height = main_paned_window.winfo_height()
        
        # 2. 하단 패널(변수/테이블)이 최소한 이만큼 보이도록 높이를 지정 (픽셀 단위)
        desired_bottom_height = 380 # [수정 2] 400 -> 380
        
        # 3. 상단 그래프 영역의 최소 높이도 지정 (너무 줄어들지 않게)
        min_top_height = 300 

        # 4. 창의 전체 높이가 두 영역의 최소 합보다 클 때만 실행
        if total_height > (desired_bottom_height + min_top_height):
            # 5. SASH의 Y좌표를 (전체 높이 - 하단 높이)로 계산
            sash_y_position = total_height - desired_bottom_height
            main_paned_window.sash_place(0, 0, sash_y_position)
        else:
            # 6. 창이 너무 작아서 최소 높이를 확보할 수 없는 경우 (예: 600px)
            #    그냥 65% 비율로 설정 (이전 코드)
            main_paned_window.sash_place(0, 0, int(total_height * 0.4))
            
    except tk.TclError: 
        pass
main_paned_window.bind("<Configure>", adjust_sash)

left_pane = ttk.Frame(top_frame, width=400)
left_pane.pack(side=tk.LEFT, fill=tk.BOTH, expand=False, padx=5, pady=5)
right_paned_window = ttk.PanedWindow(top_frame, orient=tk.HORIZONTAL)
right_paned_window.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=5, pady=5)

# --- Left Control Panel ---
left_pane.rowconfigure(1, weight=1); left_pane.columnconfigure(0, weight=1)

tree_label_frame = ttk.LabelFrame(left_pane, text="Sample / File Name")
tree_label_frame.grid(row=1, column=0, columnspan=2, sticky='nsew', pady=5)
tree_label_frame.rowconfigure(1, weight=1)
tree_label_frame.columnconfigure(0, weight=1)
tree_button_frame = ttk.Frame(tree_label_frame)
tree_button_frame.grid(row=0, column=0, sticky='ew')
ttk.Button(tree_button_frame, text="Expand All", command=expand_all_folders).pack(side=tk.LEFT, expand=True, fill=tk.X)
ttk.Button(tree_button_frame, text="Collapse All", command=collapse_all_folders).pack(side=tk.LEFT, expand=True, fill=tk.X)
ttk.Button(tree_button_frame, text="Best/Folder", command=filter_best_pce_per_folder).pack(side=tk.LEFT, expand=True, fill=tk.X)
tree_view_frame = ttk.Frame(tree_label_frame)
tree_view_frame.grid(row=1, column=0, sticky='nsew')
file_tree_scrollbar = ttk.Scrollbar(tree_view_frame)
file_tree_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
file_tree = ttk.Treeview(tree_view_frame, yscrollcommand=file_tree_scrollbar.set, selectmode='extended')
file_tree.pack(fill=tk.BOTH, expand=True)
file_tree_scrollbar.config(command=file_tree.yview)
file_tree.heading('#0', text='Sample / File', anchor='w')

jv_add_button = ttk.Button(left_pane, text="Add File(s) to J-V Plot", command=add_selected_to_jv_graph)
jv_add_button.grid(row=2, column=0, sticky='ew', padx=2)
dist_add_button = ttk.Button(left_pane, text="Add Folder(s) to Dist. Plot", command=add_selected_folder_to_dist_plot)
dist_add_button.grid(row=2, column=1, sticky='ew', padx=2)

merge_button = ttk.Button(left_pane, text="Merge Selected Folders", command=merge_selected_folders)
merge_button.grid(row=3, column=0, columnspan=2, sticky='ew', padx=5, pady=(5, 0))

manage_vars_button = ttk.Button(left_pane, text="Manage Variables...", command=open_variable_manager)
manage_vars_button.grid(row=4, column=0, columnspan=2, sticky='ew', padx=5, pady=(5, 0))

filter_frame = ttk.LabelFrame(left_pane, text="Data Filter")
filter_frame.grid(row=5, column=0, columnspan=2, sticky='ew', padx=5, pady=5)
filter_voc_min, filter_voc_max = tk.StringVar(), tk.StringVar(); filter_jsc_min, filter_jsc_max = tk.StringVar(), tk.StringVar()
filter_ff_min, filter_ff_max = tk.StringVar(), tk.StringVar(); filter_pce_min, filter_pce_max = tk.StringVar(), tk.StringVar()
filter_vars = [("Voc", filter_voc_min, filter_voc_max), ("Jsc", filter_jsc_min, filter_jsc_max),
               ("FF", filter_ff_min, filter_ff_max), ("PCE", filter_pce_min, filter_pce_max)]
for i, (name, min_var, max_var) in enumerate(filter_vars):
    ttk.Label(filter_frame, text=name).grid(row=i, column=0, padx=2)
    ttk.Entry(filter_frame, textvariable=min_var, width=8).grid(row=i, column=1)
    ttk.Label(filter_frame, text="~").grid(row=i, column=2)
    ttk.Entry(filter_frame, textvariable=max_var, width=8).grid(row=i, column=3)
filter_btn_frame = ttk.Frame(filter_frame)
filter_btn_frame.grid(row=0, column=4, rowspan=4, padx=5)
ttk.Button(filter_btn_frame, text="Apply\nFilter", command=apply_filter).pack(fill=tk.X, pady=2)
ttk.Button(filter_btn_frame, text="Reset\nFilter", command=reset_filter).pack(fill=tk.X, pady=2)

# --- Right Graph Panels ---
jv_graph_frame = ttk.LabelFrame(right_paned_window, text="J-V Curve")
dist_graph_frame = ttk.LabelFrame(right_paned_window, text="Parameter Distribution")
right_paned_window.add(jv_graph_frame, weight=1)
right_paned_window.add(dist_graph_frame, weight=1)

# --- J-V Graph Widgets ---
fig_jv, ax_jv = plt.subplots(facecolor='white', figsize=(4, 4))
jv_canvas_widget = FigureCanvasTkAgg(fig_jv, master=jv_graph_frame)
jv_canvas_widget.get_tk_widget().pack(side=tk.TOP, fill=tk.BOTH, expand=True)
jv_toolbar_frame = ttk.Frame(jv_graph_frame)
jv_toolbar_frame.pack(side=tk.BOTTOM, fill=tk.X)
toolbar = NavigationToolbar2Tk(jv_canvas_widget, jv_toolbar_frame)
toolbar.update()
ttk.Button(jv_toolbar_frame, text="Copy Graph", command=lambda: copy_figure_to_clipboard(fig_jv)).pack(side=tk.RIGHT, padx=5)

# --- [추가] Raw Table 복사 버튼 ---
ttk.Button(jv_toolbar_frame, text="Copy Raw Table", command=copy_jv_raw_data).pack(side=tk.RIGHT, padx=5)
# --- [추가 끝] ---

scan_filter_jv_var = tk.StringVar(value='All')
ttk.Label(jv_toolbar_frame, text="Scan:").pack(side=tk.RIGHT, padx=(5,2))
scan_combo_jv = ttk.Combobox(jv_toolbar_frame, textvariable=scan_filter_jv_var, values=['All', 'Reverse', 'Forward'], state='readonly', width=10)
scan_combo_jv.pack(side=tk.RIGHT, padx=2)
scan_combo_jv.bind('<<ComboboxSelected>>', lambda e: redraw_jv_graphs() if plotted_jv_items else None)
jv_control_frame = ttk.Frame(jv_graph_frame)
jv_control_frame.pack(side=tk.BOTTOM, fill=tk.X, pady=4)
plotted_list_frame = ttk.LabelFrame(jv_control_frame, text="Plotted J-V Curves (Double-click to edit)")
plotted_list_frame.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0,3))
plotted_list_tree = ttk.Treeview(plotted_list_frame, show='tree', height=4)
plotted_list_tree.pack(fill=tk.X, expand=True)
plotted_list_tree.bind("<Double-1>", lambda event: on_plotted_item_double_click(event, plotted_jv_items, plotted_list_tree))
jv_buttons_frame = ttk.Frame(plotted_list_frame)
jv_buttons_frame.pack(fill=tk.X, pady=2)
ttk.Button(jv_buttons_frame, text="Remove Selected", command=remove_selected_from_jv_graph).pack(side=tk.LEFT, expand=True, fill=tk.X, padx=2)
ttk.Button(jv_buttons_frame, text="Clear All", command=clear_jv_plot).pack(side=tk.LEFT, expand=True, fill=tk.X, padx=2)
ttk.Button(jv_buttons_frame, text="Plot Best PCEs", command=plot_best_pces).pack(side=tk.LEFT, expand=True, fill=tk.X, padx=2)
jv_axis_manage_frame = ttk.LabelFrame(jv_control_frame, text="J-V Axis Range")
jv_axis_manage_frame.pack(side=tk.LEFT)
jv_x_min_var, jv_x_max_var = tk.StringVar(), tk.StringVar()
jv_y_min_var, jv_y_max_var = tk.StringVar(), tk.StringVar()
ttk.Label(jv_axis_manage_frame, text="X min:").grid(row=0, column=0); ttk.Entry(jv_axis_manage_frame, textvariable=jv_x_min_var, width=7).grid(row=0, column=1)
ttk.Label(jv_axis_manage_frame, text="X max:").grid(row=0, column=2); ttk.Entry(jv_axis_manage_frame, textvariable=jv_x_max_var, width=7).grid(row=0, column=3)
ttk.Label(jv_axis_manage_frame, text="Y min:").grid(row=1, column=0); ttk.Entry(jv_axis_manage_frame, textvariable=jv_y_min_var, width=7).grid(row=1, column=1)
ttk.Label(jv_axis_manage_frame, text="Y max:").grid(row=1, column=2); ttk.Entry(jv_axis_manage_frame, textvariable=jv_y_max_var, width=7).grid(row=1, column=3)
ttk.Button(jv_axis_manage_frame, text="Apply", command=update_jv_axis_limits).grid(row=2, column=1, rowspan=2, padx=4, sticky='ns')
ttk.Button(jv_axis_manage_frame, text="Auto", command=lambda: auto_scale_jv_axes(smart_scale=True)).grid(row=3, column=2, rowspan=2, padx=4, sticky='ns')

# --- Distribution Graph Widgets ---
# [수정] 2x2/1x4 (컨트롤 하단) 및 4x1 (컨트롤 우측)을 지원하는 동적 레이아웃으로 변경

# 1. (신규) 왼쪽(그래프)과 오른쪽(컨트롤)을 나누는 프레임 2개 생성
dist_plot_container_frame = ttk.Frame(dist_graph_frame)
dist_controls_area_frame = ttk.Frame(dist_graph_frame) # [중요] 모든 컨트롤 위젯의 부모 프레임


# 2. 왼쪽 프레임(dist_plot_container_frame)에 그래프/툴바 배치
fig_dist, axs_dist = plt.subplots(2, 2, facecolor='white', figsize=(7, 4.5))
    
dist_canvas_widget = FigureCanvasTkAgg(fig_dist, master=dist_plot_container_frame) # 부모: 왼쪽 프레임
dist_canvas_widget.get_tk_widget().pack(side=tk.TOP, fill=tk.BOTH, expand=True)
    
dist_toolbar_frame = ttk.Frame(dist_plot_container_frame) # 부모: 왼쪽 프레임
dist_toolbar_frame.pack(side=tk.BOTTOM, fill=tk.X)
    
dist_toolbar = NavigationToolbar2Tk(dist_canvas_widget, dist_toolbar_frame)
dist_toolbar.update()
    
# 3. 툴바 프레임(dist_toolbar_frame)에 커스텀 버튼들 추가
ttk.Button(dist_toolbar_frame, text="Copy All", command=lambda: copy_figure_to_clipboard(fig_dist)).pack(side=tk.RIGHT, padx=5)

dist_layout_var = tk.StringVar(value='2x2')
ttk.Label(dist_toolbar_frame, text="Layout:").pack(side=tk.RIGHT, padx=(5,2))
layout_combo = ttk.Combobox(dist_toolbar_frame, textvariable=dist_layout_var, values=['2x2', '1x4', '4x1'], state='readonly', width=6)
layout_combo.pack(side=tk.RIGHT, padx=2)
layout_combo.bind('<<ComboboxSelected>>', lambda e: change_dist_layout())

scan_filter_dist_var = tk.StringVar(value='All')
ttk.Label(dist_toolbar_frame, text="Scan:").pack(side=tk.RIGHT, padx=(5,2))
scan_combo_dist = ttk.Combobox(dist_toolbar_frame, textvariable=scan_filter_dist_var, values=['All', 'Reverse', 'Forward'], state='readonly', width=10)
scan_combo_dist.pack(side=tk.RIGHT, padx=2)
scan_combo_dist.bind('<<ComboboxSelected>>', lambda e: redraw_dist_plot())


# 4. (신규) 오른쪽 프레임(dist_controls_area_frame)에 모든 컨트롤을 배치
#    [중요] 3개의 자식 프레임을 global로 선언
global dist_list_frame, dist_axis_manage_frame, dist_export_frame

# 4-1. Plotted Distributions 리스트
dist_list_frame = ttk.LabelFrame(dist_controls_area_frame, text="Plotted Distributions (Double-click to edit)")
# dist_list_frame.pack() # <- [삭제] 5번 섹션에서 초기 배치

dist_list_tree = ttk.Treeview(dist_list_frame, show='tree', height=4)
dist_list_tree.pack(fill=tk.BOTH, expand=True, padx=2, pady=2) # fill=BOTH, expand=True로 변경
    
dist_buttons_frame = ttk.Frame(dist_list_frame)
dist_buttons_frame.pack(fill=tk.X, pady=2)
ttk.Button(dist_buttons_frame, text="Remove Selected", command=remove_selected_from_dist_graph).pack(side=tk.LEFT, expand=True, fill=tk.X, padx=2)
ttk.Button(dist_buttons_frame, text="Clear All", command=clear_dist_plot).pack(side=tk.LEFT, expand=True, fill=tk.X, padx=2)
    
# 4-2. Y-Axis Range
dist_axis_manage_frame = ttk.LabelFrame(dist_controls_area_frame, text="Y-Axis Range")
# dist_axis_manage_frame.pack() # <- [삭제] 5번 섹션에서 초기 배치

dist_voc_min_var, dist_voc_max_var = tk.StringVar(), tk.StringVar()
dist_jsc_min_var, dist_jsc_max_var = tk.StringVar(), tk.StringVar()
dist_ff_min_var, dist_ff_max_var = tk.StringVar(), tk.StringVar()
dist_pce_min_var, dist_pce_max_var = tk.StringVar(), tk.StringVar()
params = ["Voc", "Jsc", "FF", "PCE"]
vars_grid = [(dist_voc_min_var, dist_voc_max_var), (dist_jsc_min_var, dist_jsc_max_var),
             (dist_ff_min_var, dist_ff_max_var), (dist_pce_min_var, dist_pce_max_var)]
for i, param in enumerate(params):
    ttk.Label(dist_axis_manage_frame, text=f"{param} min:").grid(row=i, column=0, sticky='e', padx=(5,0)) # 여백 추가
    ttk.Entry(dist_axis_manage_frame, textvariable=vars_grid[i][0], width=7).grid(row=i, column=1)
    ttk.Label(dist_axis_manage_frame, text="max:").grid(row=i, column=2, sticky='e')
    ttk.Entry(dist_axis_manage_frame, textvariable=vars_grid[i][1], width=7).grid(row=i, column=3)
        
dist_axis_buttons_frame = ttk.Frame(dist_axis_manage_frame)
dist_axis_buttons_frame.grid(row=0, column=4, rowspan=4, padx=5, pady=4) # 여백 추가
ttk.Button(dist_axis_buttons_frame, text="Apply", command=update_dist_axis_limits).pack(expand=True, fill=tk.BOTH)
ttk.Button(dist_axis_buttons_frame, text="Auto", command=auto_scale_dist_axes).pack(expand=True, fill=tk.BOTH)

# 4-3. Copy Subplots
dist_export_frame = ttk.LabelFrame(dist_controls_area_frame, text="Copy Subplots")
# dist_export_frame.pack() # <- [삭제] 5번 섹션에서 초기 배치

ttk.Button(dist_export_frame, text="Copy Voc...", command=lambda: copy_subplot_to_clipboard(0)).pack(padx=5, pady=2, fill=tk.X)
ttk.Button(dist_export_frame, text="Copy Jsc...", command=lambda: copy_subplot_to_clipboard(1)).pack(padx=5, pady=2, fill=tk.X)
ttk.Button(dist_export_frame, text="Copy FF...", command=lambda: copy_subplot_to_clipboard(2)).pack(padx=5, pady=2, fill=tk.X)
ttk.Button(dist_export_frame, text="Copy PCE...", command=lambda: copy_subplot_to_clipboard(3)).pack(padx=5, pady=2, fill=tk.X)


# 5. [중요] 초기 레이아웃(2x2) 설정
# 5-1. 메인 프레임 배치: 그래프(TOP), 컨트롤(BOTTOM)
dist_plot_container_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=5, pady=5)
dist_controls_area_frame.pack(side=tk.BOTTOM, fill=tk.X, expand=False, padx=5, pady=5)

# 5-2. 컨트롤 프레임 내부 자식 3개 배치: 가로 3분할 (LEFT, LEFT, LEFT)
dist_list_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, pady=0, padx=(0,5))
dist_axis_manage_frame.pack(side=tk.LEFT, fill=tk.Y, expand=False, pady=0, padx=5)
dist_export_frame.pack(side=tk.LEFT, fill=tk.Y, expand=False, pady=0, padx=(5,0))

# --- Bottom Results Tabs ---
bottom_paned_window = ttk.PanedWindow(bottom_frame, orient=tk.HORIZONTAL)
bottom_paned_window.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)

var_manage_frame = ttk.LabelFrame(bottom_paned_window, text="Experimental Variables")
bottom_paned_window.add(var_manage_frame, weight=1) 

var_header_frame = ttk.Frame(var_manage_frame)
var_header_frame.pack(fill='x', padx=10, pady=5)
ttk.Label(var_header_frame, text="Current Sample:").pack(side=tk.LEFT)
current_sample_label = tk.StringVar(value="[No sample selected]")
ttk.Label(var_header_frame, textvariable=current_sample_label, font="-weight bold").pack(side=tk.LEFT, padx=5)

var_grid_frame = ttk.Frame(var_manage_frame)
var_grid_frame.pack(fill='x', padx=10, pady=5)

var_entry_widgets.clear()
var_checkbox_vars.clear()

for i, col in enumerate(variable_columns):
    row, col_idx = i % 5, (i // 5) * 3

    chk_var = tk.BooleanVar()
    chk = ttk.Checkbutton(var_grid_frame, text=f"{col}:", variable=chk_var)
    chk.grid(row=row, column=col_idx, sticky='w', padx=5, pady=2)

    var = tk.StringVar()
    entry = ttk.Entry(var_grid_frame, textvariable=var, width=15)
    entry.grid(row=row, column=col_idx + 1, sticky='w', padx=5, pady=2)

    var_entry_widgets[col] = (var, entry)
    var_checkbox_vars[col] = chk_var

    if col == "Perovskite":
        # [신규] "Build..."와 "..." 버튼 2개를 나란히 담을 프레임 생성
        button_frame = ttk.Frame(var_grid_frame)
        
        # [수정] 프레임을 기존 버튼 위치(col_idx + 2)에 배치합니다.
        # padx=0으로 설정하고, 개별 버튼에 padx를 줍니다.
        button_frame.grid(row=row, column=col_idx + 2, sticky='w', padx=0)

        # [수정] "Build..." 버튼의 부모를 'var_grid_frame' -> 'button_frame'으로 변경
        # .grid() 대신 .pack(side=tk.LEFT) 사용
        ttk.Button(
            button_frame, 
            text="Build...", 
            command=lambda s=current_sample_label: PerovskiteBuilderWindow(root, s.get()), 
            width=5
        ).pack(side=tk.LEFT, padx=(2, 2)) # 왼쪽에 배치
        
        entry.config(state='readonly')

        # [신규] "..." (상세정보) 버튼을 'button_frame'에 추가
        btn = ttk.Button(
            button_frame, # 부모를 'button_frame'으로
            text="...",
            width=3,
            command=lambda c=col, v=var: open_variable_selector(c, v)
        )
        btn.pack(side=tk.LEFT, padx=(0, 2)) # "Build..." 버튼 오른쪽에 배치

    elif col in VARIABLE_PRESETS:
        # (이 부분은 기존 코드와 동일합니다)
        btn = ttk.Button(
            var_grid_frame,
            text="...",
            width=3,
            command=lambda c=col, v=var: open_variable_selector(c, v)
        )
        btn.grid(row=row, column=col_idx + 2, sticky='w', padx=2)
        entry.config(state='readonly')
    elif col in VARIABLE_PRESETS:
        btn = ttk.Button(
            var_grid_frame,
            text="...",
            width=3,
            command=lambda c=col, v=var: open_variable_selector(c, v)
        )
        btn.grid(row=row, column=col_idx + 2, sticky='w', padx=2)
        entry.config(state='readonly')

# "Contact" 항목 (i=6)의 위치 계산:
# row = 6 % 5 = 1
# col_idx = (6 // 5) * 3 = 3
# Contact의 Entry 위젯은 column = col_idx + 1 = 4 에 있습니다.
# 따라서 버튼은 row=2, column=4 부터 시작합니다.

save_vars_button = ttk.Button(
    var_grid_frame,  # 1. 부모를 var_grid_frame으로 변경
    text="Save Variables for Selected Sample", 
    command=lambda: save_variables_for_selected_sample()
)
# 2. .pack() 대신 .grid() 사용
save_vars_button.grid(
    row=2,           # Contact의 다음 행(row)
    column=4,        # Contact의 Entry와 같은 열(column)
    columnspan=2,    # Entry와 "..." 버튼의 너비를 합친 만큼(2칸) 차지
    sticky='ew',     # 좌우로 꽉 채우기
    padx=5, 
    pady=(10, 2)     # 위쪽 여백 10, 아래쪽 2
)

view_details_button = ttk.Button(
    var_grid_frame,  # 1. 부모를 var_grid_frame으로 변경
    text="View Process Details",
    command=view_process_details
)
# 2. .pack() 대신 .grid() 사용
view_details_button.grid(
    row=3,           # "Save" 버튼의 다음 행(row)
    column=4,        # "Save" 버튼과 같은 열(column)
    columnspan=2,    # 2칸 차지
    sticky='ew',     # 좌우로 꽉 채우기
    padx=5, 
    pady=2
)

# 이제 var_manage_frame은 비어있게 되므로, 
# 이 프레임을 생성하고 .pack()하는 코드를 제거하거나 주석 처리할 수 있습니다.
notebook = ttk.Notebook(bottom_paned_window)
bottom_paned_window.add(notebook, weight=2) 

tab1 = ttk.Frame(notebook); tab2 = ttk.Frame(notebook); tab3 = ttk.Frame(notebook)
notebook.add(tab1, text="All Devices"); notebook.add(tab2, text="Best Device by Sample"); notebook.add(tab3, text="Statistics by Sample")

# --- Tab 1: All Devices Table ---
all_devices_frame = ttk.Frame(tab1)
all_devices_frame.pack(fill=tk.BOTH, expand=True)
copy_btn1 = ttk.Button(all_devices_frame, text="Copy Table", command=lambda: copy_treeview_to_clipboard(pce_table_all))
copy_btn1.pack(anchor='ne', padx=5, pady=2)
pce_table_all_frame = ttk.Frame(all_devices_frame)
pce_table_all_frame.pack(fill=tk.BOTH, expand=True)
pce_table_all = ttk.Treeview(pce_table_all_frame, columns=('File', 'Scan', 'Voc (V)', 'Jsc (mA/cm2)', 'FF (%)', 'PCE (%)', 'Rs (Ω·cm²)','Rsh (Ω·cm²)'))
pce_table_all.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
v_scroll1 = ttk.Scrollbar(pce_table_all_frame, orient=tk.VERTICAL, command=pce_table_all.yview)
h_scroll1 = ttk.Scrollbar(all_devices_frame, orient=tk.HORIZONTAL, command=pce_table_all.xview)
pce_table_all.configure(yscrollcommand=v_scroll1.set, xscrollcommand=h_scroll1.set)
v_scroll1.pack(side=tk.RIGHT, fill=tk.Y); h_scroll1.pack(side=tk.BOTTOM, fill=tk.X)
pce_table_all.column('#0', width=0, stretch=tk.NO); pce_table_all.heading('#0', text='', anchor='w')
for col in pce_table_all['columns']:
    pce_table_all.heading(col, text=col, anchor='center'); pce_table_all.column(col, anchor='center', width=100)
pce_table_all.column('File', width=300, anchor='w')
pce_table_all.column('Scan', width=80)

# --- Tab 2: Best Devices Table ---
best_devices_frame = ttk.Frame(tab2)
best_devices_frame.pack(fill=tk.BOTH, expand=True)
copy_btn2 = ttk.Button(best_devices_frame, text="Copy Table", command=lambda: copy_treeview_to_clipboard(pce_table_best))
copy_btn2.pack(anchor='ne', padx=5, pady=2)
pce_table_best_frame = ttk.Frame(best_devices_frame)
pce_table_best_frame.pack(fill=tk.BOTH, expand=True)
pce_table_best = ttk.Treeview(pce_table_best_frame, columns=('Sample', 'File', 'Scan', 'Voc (V)', 'Jsc (mA/cm2)', 'FF (%)', 'PCE (%)', 'Rs (Ω·cm²)', 'Rsh (Ω·cm²)'))
pce_table_best.pack(fill=tk.BOTH, expand=True)
for col in pce_table_best['columns']:
    pce_table_best.heading(col, text=col, anchor='center'); pce_table_best.column(col, anchor='center', width=120)
pce_table_best.column('#0', width=0, stretch=tk.NO)
pce_table_best.column('Scan', width=80)

# --- Tab 3: Statistics Table ---
stats_frame = ttk.Frame(tab3)
stats_frame.pack(fill=tk.BOTH, expand=True)
copy_btn3 = ttk.Button(stats_frame, text="Copy Table", command=lambda: copy_treeview_to_clipboard(pce_table_stats))
copy_btn3.pack(anchor='ne', padx=5, pady=2)
pce_table_stats_frame = ttk.Frame(stats_frame)
pce_table_stats_frame.pack(fill=tk.BOTH, expand=True)
pce_table_stats = ttk.Treeview(stats_frame, columns=('Sample', 'Count', 'Voc (V)', 'Jsc (mA/cm2)', 'FF (%)', 'PCE (%)', 'Rs (Ω·cm²)', 'Rsh (Ω·cm²)'))
pce_table_stats.pack(fill=tk.BOTH, expand=True)
for col in pce_table_stats['columns']:
    pce_table_stats.heading(col, text=col, anchor='center'); pce_table_stats.column(col, anchor='center', width=150)
pce_table_stats.column('#0', width=0, stretch=tk.NO)

# --- Signature / Copyright Label ---
signature_frame = ttk.Frame(bottom_frame)
signature_frame.pack(side=tk.BOTTOM, fill=tk.X, padx=10, pady=(0, 5))
signature_text = "Solar Cell Analysis Tool | Developed by Hyoungwoo Kwon"
ttk.Label(signature_frame, text=signature_text, font=('Helvetica', 9), anchor='e').pack(fill=tk.X)

# --- Functions to connect UI elements ---
def on_file_tree_select(event):
    """파일 트리에서 샘플 폴더를 선택하면 변수 관리자 UI를 업데이트합니다."""
    selected_items = file_tree.selection()
    if not selected_items:
        current_sample_label.set("[No sample selected]")
        for col in variable_columns:
            var_entry_widgets[col][0].set("")
            # 'readonly' 상태인 엔트리도 비활성화합니다.
            if col in VARIABLE_PRESETS or col == "Perovskite":
                var_entry_widgets[col][1].config(state='disabled')
            else:
                var_entry_widgets[col][1].config(state='disabled')

            if col in var_checkbox_vars:
                var_checkbox_vars[col].set(False)
        save_vars_button.config(state='disabled')
        return

    selected_item_id = selected_items[0]
    if file_tree.parent(selected_item_id):
        selected_item_id = file_tree.parent(selected_item_id)

    sample_name = file_tree.item(selected_item_id, 'text')
    current_sample_label.set(sample_name)

    sample_vars = experimental_variables.get(sample_name, {})
    for col in variable_columns:
        var_entry_widgets[col][0].set(sample_vars.get(col, ""))

        # 'readonly'가 아닌 일반 엔트리만 'normal'로 설정합니다.
        if col not in VARIABLE_PRESETS and col != "Perovskite":
            var_entry_widgets[col][1].config(state='normal')
        else:
            # 'readonly' 항목은 다시 활성화합니다.
            var_entry_widgets[col][1].config(state='readonly')

        if col in var_checkbox_vars:
            var_checkbox_vars[col].set(sample_vars.get(f"{col}_is_var", False))

    save_vars_button.config(state='normal')

def save_variables_for_selected_sample():
    """변수 관리자 패널의 현재 값을 선택된 샘플에 저장합니다."""
    global current_display_df, experimental_variables, process_details # <-- 1. process_details를 global로 가져옵니다.
    sample_name = current_sample_label.get()
    if sample_name == "[No sample selected]":
        messagebox.showwarning("No Sample", "Please select a sample folder from the list first.")
        return

    if sample_name not in experimental_variables:
        experimental_variables[sample_name] = {}

    control_vars_to_fill = {}
    for col in variable_columns:
        value = var_entry_widgets[col][0].get()
        is_variable = var_checkbox_vars[col].get() if col in var_checkbox_vars else False

        experimental_variables[sample_name][col] = value
        experimental_variables[sample_name][f"{col}_is_var"] = is_variable

        if not is_variable and value:
            control_vars_to_fill[col] = value

    if control_vars_to_fill:
        msg = "The following control variables (unchecked) will be applied to ALL other samples:\n\n"
        msg += "\n".join([f"{col}: {val}" for col, val in control_vars_to_fill.items()])
        msg += "\n\nProceed?"
        if messagebox.askyesno("Apply Control Variables?", msg, parent=root):
            all_samples = current_display_df['Sample'].unique()
            for s in all_samples:
                if s == sample_name: 
                    continue
                if s not in experimental_variables:
                    experimental_variables[s] = {}
                
                for col, val in control_vars_to_fill.items():
                    is_var_key = f"{col}_is_var"
                    if is_var_key not in experimental_variables[s]:
                        experimental_variables[s][is_var_key] = False
                    
                    # 1. 'experimental_variables'의 값을 복사합니다. (기존 로직)
                    experimental_variables[s][col] = val

                    # --- NEW: 'process_details'도 함께 복사합니다. (추가된 로직) ---
                    # 'val' (예: "NiOx + PEDOT")을 개별 재료로 분리합니다.
                    materials_to_copy = [m.strip() for m in val.split(' + ') if m.strip()]
                    
                    for material in materials_to_copy:
                        # 원본 샘플(현재 선택된 샘플)에서 공정 상세정보를 가져옵니다.
                        source_details = process_details.get(sample_name, {}).get(col, {}).get(material, {})
                        
                        if source_details: # 복사할 상세정보가 있다면
                            # 대상 샘플(s)의 딕셔너리 구조를 생성합니다.
                            if s not in process_details:
                                process_details[s] = {}
                            if col not in process_details[s]:
                                process_details[s][col] = {}
                                
                            # 대상 샘플(s)에 상세정보를 복사합니다.
                            process_details[s][col][material] = source_details.copy()
                    # --- END NEW ---

    merged_vars_df = pd.DataFrame.from_dict(
        experimental_variables, 
        orient='index'
    ).reset_index().rename(columns={'index': 'Sample'})

    all_var_cols_to_consider = variable_columns + [f"{col}_is_var" for col in variable_columns]
    cols_to_drop = [
        col for col in all_var_cols_to_consider 
        if col in current_display_df.columns and col != 'Sample'
    ]

    df_without_vars = current_display_df.drop(columns=cols_to_drop, errors='ignore')
    current_display_df = pd.merge(df_without_vars, merged_vars_df, on="Sample", how="left")

    refresh_all_views(current_display_df)
    messagebox.showinfo("Variables Saved", f"Experimental variables for '{sample_name}' have been saved.")

# Bind the function to the file tree selection
file_tree.bind('<<TreeviewSelect>>', on_file_tree_select)

# --- [추가] 마우스 오른쪽 클릭 이벤트를 바인딩합니다. ---
file_tree.bind("<Button-3>", on_file_tree_right_click)
# --- [추가 끝] ---


# Initially disable variable entries
for col in variable_columns:
    if col in VARIABLE_PRESETS or col == "Perovskite":
         var_entry_widgets[col][1].config(state='disabled')
    else:
         var_entry_widgets[col][1].config(state='disabled')
save_vars_button.config(state='disabled')


# --- Initial Draw ---
redraw_jv_graphs()
redraw_dist_plot()
root.mainloop()

